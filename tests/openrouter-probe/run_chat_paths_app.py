"""run_chat_paths_app.py -- Comprehensive all-layers chat probe for GLOW.

Tests every distinct code path in the GLOW chat feature across 9 layers:

LAYER 0 -- Document ingestion (_load_document_text equivalent)
  Path 00a  .md direct read_text()
  Path 00b  .docx MarkItDown extraction
  Path 00c  .xlsx MarkItDown extraction
  Path 00d  .pptx MarkItDown extraction
  Path 00e  .pdf (text-based) MarkItDown extraction
  Path 00f  .epub MarkItDown extraction
  Path 00g  .html MarkItDown extraction
  Path 00h  .csv MarkItDown extraction
  Path 00i  _build_document_context long-doc truncation (>8000 chars)
  Path 00j  _build_document_context short-doc passthrough
  Path 00k  _build_document_context keyword snippet extraction

LAYER 1 -- Raw OpenRouter gateway (ai_gateway.py)
  Path 01  Free model, ACB font question
  Path 02  Paid fallback model (gpt-4o-mini) echo
  Path 03  gateway chat() free-first
  Path 04  gateway chat() with document context
  Path 05  gateway chat() multi-turn conversation history
  Path 06  gateway chat() rich ACB audit question
  Path 07  gateway chat() uncertainty trigger
  Path 08  quota/budget status probe (no model call)

LAYER 2 -- Document tools (ToolRegistry)
  Path 09  get_document_stats
  Path 10  list_headings
  Path 11  check_heading_hierarchy (valid)
  Path 12  check_heading_hierarchy (H2->H4 skip)
  Path 13  find_faux_headings (violations)
  Path 14  find_faux_headings (clean)
  Path 15  check_list_structure (deep nesting)
  Path 16  check_list_structure (clean)
  Path 17  estimate_reading_order
  Path 18  find_section (match)
  Path 19  find_section (no match)
  Path 20  search_text (match)
  Path 21  search_text (no match)
  Path 22  extract_table (index 0)
  Path 23  extract_table (not found)
  Path 24  get_images
  Path 24b get_document_summary
  Path 24c get_section_content (found)
  Path 24d get_section_content (not found)
  Path 24e get_decisions_and_actions
  Path 24f get_what_passes

LAYER 3 -- Compliance Agent tools
  Path 25  run_accessibility_audit (violations)
  Path 26  run_accessibility_audit (clean)
  Path 27  get_compliance_score
  Path 28  get_critical_findings
  Path 29  get_auto_fixable_findings

LAYER 4 -- Content Agent tools
  Path 30  check_emphasis_patterns (italic + bold abuse)
  Path 31  check_emphasis_patterns (clean)
  Path 32  check_link_text (bare URL + click here)
  Path 33  check_link_text (clean)
  Path 34  check_reading_level (high complexity)
  Path 35  check_reading_level (plain language)
  Path 36  check_alignment_hints (center found)
  Path 37  check_alignment_hints (clean)

LAYER 5 -- Remediation Agent tools
  Path 38  explain_rule (ACB-NO-ITALIC)
  Path 39  explain_rule (ACB-FAUX-HEADING)
  Path 40  explain_rule (ACB-ALIGNMENT)
  Path 41  explain_rule (ACB-TABLE-HEADERS)
  Path 42  explain_rule (unknown rule)
  Path 43  suggest_fix (ACB-FONT-SIZE-BODY)
  Path 44  prioritize_findings (heuristic)
  Path 45  estimate_fix_impact (no cache)
  Path 46  check_image_alt_text (2 of 3 missing)
  Path 47  check_image_alt_text (all present)
  Path 48  check_image_alt_text (no images)

LAYER 6 -- Pre-flight keyword dispatcher
  Path 49  audit / compliance
  Path 50  fix / remediate
  Path 51  heading / hierarchy
  Path 52  italic / bold / emphasis
  Path 53  link / url
  Path 54  image / alt text
  Path 55  list / bullet
  Path 56  reading level
  Path 57  find section
  Path 58  search term
  Path 59  table / data
  Path 60  ACB rule ID
  Path 61  no keywords default
  Path 62  combined italic + heading

LAYER 7 -- Full GLOW pipeline (system prompt + pre-flight + gateway)
  Path 63  All critical violations
  Path 64  Prioritize fixes
  Path 65  Heading structure
  Path 66  Explain italic fix
  Path 67  Table check
  Path 68  Reading level
  Path 69  Section budget
  Path 70  Multi-turn follow-up
  Path 70b Document summary
  Path 70c Section content (financial reserves)
  Path 70d Decisions and action items

LAYER 8 -- Vision and audio gateway
  Path 71  describe_image() annotated PNG with text
  Path 72  describe_image() tiny synthetic PNG
  Path 73  transcribe() silent WAV
  Path 74  transcribe() real mp3 (optional)

All responses written to tests/openrouter-probe/chat/paths/

Usage:
    cd s:\\code\\glow
    python tests\\openrouter-probe\\run_chat_paths_app.py [--layer N]

    --layer N  run only that layer (0-8)
"""

from __future__ import annotations

import argparse
import io
import json
import re
import struct
import sys
import tempfile
import textwrap
import wave
import zlib
from pathlib import Path
from typing import Any

# ---------------------------------------------------------------------------
# Path setup
# ---------------------------------------------------------------------------

PROBE_ROOT = Path(__file__).parent
REPO_ROOT = PROBE_ROOT.parents[1]
sys.path.insert(0, str(REPO_ROOT / "web" / "src"))
sys.path.insert(0, str(REPO_ROOT / "desktop" / "src"))

import requests

from acb_large_print_web.credentials import get_openrouter_api_key
from acb_large_print_web.ai_gateway import (
    chat as gateway_chat,
    describe_image,
    get_quota_status,
    transcribe,
)
from acb_large_print_web.chat_handler import (
    DocumentContext,
    ToolRegistry,
)
from acb_large_print_web.routes.chat import (
    _CHAT_SYSTEM_PROMPT,
    _build_document_context,
)

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

OPENROUTER_BASE = "https://openrouter.ai/api/v1"
OUT_DIR = PROBE_ROOT / "chat" / "paths"
OUT_DIR.mkdir(parents=True, exist_ok=True)

API_KEY = get_openrouter_api_key()
SESSION = "probe-all-layers-001"

RESULTS: list[dict[str, Any]] = []
LAYER_COUNTERS: dict[int, dict[str, int]] = {}

# ---------------------------------------------------------------------------
# Rich synthetic documents
# ---------------------------------------------------------------------------

RICH_BAD_DOC = """\
# BITS Board Meeting -- April 2026

*Minutes recorded by the Secretary*

## 1. Attendance

All board members were present. The meeting was called to order at 9:00 AM by
President Jeff Bishop.

**Thomas Hartley, Treasurer, noted that the following agenda items had been revised.**

## 2. Financial Report

### 2a. Budget Overview

| Quarter | Revenue   | Expenses  | Net       |
|---------|-----------|-----------|-----------|
| Q1 2026 | $142,000  | $118,000  | $24,000   |
| Q2 2026 | $158,000  | $131,000  | $27,000   |
| Q3 2026 | $155,000  | $127,500  | $27,500   |
| Q4 2026 | $161,000  | $134,000  | $27,000   |

The treasurer presented the annual budget summary. The organization is operating
within projected parameters. See https://internal.bits-acb.org/financials for
the full breakdown or [click here](https://internal.bits-acb.org/financials) for
the detailed report.

#### 2b. Reserves

The organization maintains a three-month operating reserve of $354,000. The
board voted unanimously to maintain this level throughout 2026.

### 2c. Grant Activity

![](grant-chart.png)

Three grants are currently active. Two new applications were submitted in March.

## 3. Technology Committee Report

**All committee members submitted written reports in advance.**

*The committee chair reviewed each report in order.*

### 3a. Website Accessibility Audit

The website was audited against WCAG 2.2 AA. Critical findings included
missing alt text on 14 images, 3 instances of color-only information, and
2 form fields without labels. Remediation is scheduled for May 2026.

<center>This summary was prepared by the Technology Committee.</center>

### 3b. Software Procurement

The board reviewed proposals from three vendors. All vendors were evaluated on
price, accessibility compliance, and long-term support commitments. The
procurement timeline extends through Q3 2026. Evaluations require substantial
cross-functional review involving procurement, technology, legal, and
accessibility specialists working in coordination over an extended review period
that may require supplemental budget authorization for external consulting.

## 4. Membership Update

- New members in Q1: 47
- Renewals: 312
- Total active membership: 1,204
  - Chapter members: 876
  - At-large members: 328
    - International at-large: 42
      - Canada: 31
        - Ontario chapter: 19

## 5. Action Items

1. Finance committee to present revised budget by May 15
2. Technology committee to complete website remediation by May 31
3. Secretary to distribute minutes within 5 business days
4. President to confirm annual meeting venue by April 30

## 6. Next Meeting

The next board meeting is scheduled for July 8, 2026 at 9:00 AM.

![Meeting agenda distributed to all board members before the session](agenda.png)
![](missing-alt.png)
"""

RICH_CLEAN_DOC = """\
# BITS Annual Newsletter -- Spring 2026

## Welcome from the President

The Blind Information Technology Solutions organization has had a strong start
to 2026. Membership is up, our website meets WCAG 2.2 AA, and our technology
committee launched three new accessible tools.

## Events

- Annual Conference: June 14-16, 2026
- Regional Meetups: See the [Events Calendar](https://bits-acb.org/events)
- Webinar Series: Registration open at [BITS Webinars](https://bits-acb.org/webinars)

## Technology News

### Large Print Toolkit

The GLOW toolkit now supports Word, Excel, PowerPoint, PDF, and ePub
documents. All outputs meet ACB Large Print Guidelines.

### Screen Reader Compatibility

All BITS web properties now pass automated NVDA and JAWS testing.

## Resources

![BITS logo -- a bold B, I, T, S in blue on a white background](bits-logo.png)

See the [BITS Resource Library](https://bits-acb.org/resources) for guides,
templates, and training materials.
"""

SKIPPED_HEADING_DOC = """\
# Annual Report

## Executive Summary

#### Detailed Findings

This heading skips from H2 to H4, violating heading hierarchy rules.

## Recommendations

Content here.
"""

PLAIN_LANGUAGE_DOC = """\
# Meeting Notes

We met on April 8.

## What we decided

We will buy new software. The cost is $5,000.

## Next steps

- Order the software by May 1.
- Train staff by June 1.
"""

LONG_DOC = (RICH_BAD_DOC * 6)[:12000]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _auth_headers() -> dict[str, str]:
    return {
        "Authorization": f"Bearer {API_KEY}",
        "HTTP-Referer": "https://glow.bits-acb.org",
        "X-Title": "GLOW All-Layers Probe",
        "Content-Type": "application/json",
    }


def _chat_raw(model: str, messages: list[dict], max_tokens: int = 300) -> dict:
    resp = requests.post(
        f"{OPENROUTER_BASE}/chat/completions",
        headers=_auth_headers(),
        json={"model": model, "messages": messages,
              "max_tokens": max_tokens, "temperature": 0.0},
        timeout=90,
    )
    resp.raise_for_status()
    return resp.json()


def _record(path_id: str, name: str, ok: bool, answer: str, raw: object,
            layer: int, notes: str = "") -> None:
    entry: dict[str, Any] = {
        "path_id": path_id, "layer": layer, "name": name,
        "ok": ok, "answer_preview": answer[:600], "notes": notes,
    }
    RESULTS.append(entry)
    if layer not in LAYER_COUNTERS:
        LAYER_COUNTERS[layer] = {"pass": 0, "fail": 0}
    LAYER_COUNTERS[layer]["pass" if ok else "fail"] += 1
    slug = f"path_{path_id}"
    (OUT_DIR / f"{slug}_answer.txt").write_text(answer, encoding="utf-8")
    (OUT_DIR / f"{slug}_raw.json").write_text(
        json.dumps(raw, indent=2, default=str), encoding="utf-8")
    status = "PASS" if ok else "FAIL"
    print(f"  [{status}] {path_id}: {name}")
    if notes:
        print(f"          {notes}")


def _ctx(text: str, filename: str = "test.md") -> tuple[DocumentContext, ToolRegistry]:
    ctx = DocumentContext(text, filename)
    return ctx, ToolRegistry(ctx)


def _full_pipeline(question: str, doc: str = RICH_BAD_DOC) -> tuple[str, bool, list[str]]:
    """Simulate routes/chat.py: dispatch -> GLOW system prompt -> gateway."""
    _, tools = _ctx(doc)
    preflight = tools.dispatch_for_question(question)
    full_system = f"{_CHAT_SYSTEM_PROMPT}\n\n{preflight}"
    doc_context = _build_document_context(question, doc)
    answer, escalated = gateway_chat(
        question=question,
        system_prompt=full_system,
        session_hash=SESSION,
        document_text=doc_context,
    )
    tools_run = [
        ln.strip("[]")
        for ln in preflight.splitlines()
        if ln.startswith("[") and "===" not in ln
    ]
    return answer, escalated, tools_run


# ---------------------------------------------------------------------------
# Synthetic file builders
# ---------------------------------------------------------------------------

def _make_png_with_text(text: str = "ACB Large Print") -> bytes:
    try:
        from PIL import Image, ImageDraw
        img = Image.new("RGB", (400, 100), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.text((10, 35), text, fill=(0, 0, 0))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except ImportError:
        return _make_tiny_png()


def _make_tiny_png() -> bytes:
    def _chunk(tag: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)
    ihdr = _chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    idat = _chunk(b"IDAT", zlib.compress(b"\x00\xFF\x00\x00"))
    iend = _chunk(b"IEND", b"")
    return b"\x89PNG\r\n\x1a\n" + ihdr + idat + iend


def _make_silent_wav(path: Path, duration: float = 1.0, rate: int = 16000) -> Path:
    buf = io.BytesIO()
    with wave.open(buf, "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(rate)
        wf.writeframes(b"\x00\x00" * int(rate * duration))
    path.write_bytes(buf.getvalue())
    return path


def _make_docx(path: Path) -> Path:
    from docx import Document
    doc = Document()
    doc.add_heading("BITS Board Meeting -- April 2026", level=1)
    doc.add_heading("1. Financial Report", level=2)
    p = doc.add_paragraph()
    run = p.add_run("This text is italic and violates ACB-NO-ITALIC. ")
    run.italic = True
    p2 = doc.add_paragraph()
    run2 = p2.add_run(
        "All staff must read this notice. The notice applies to all departments.")
    run2.bold = True
    doc.add_heading("2. Budget", level=2)
    table = doc.add_table(rows=3, cols=3)
    hdr = table.rows[0].cells
    hdr[0].text = "Quarter"
    hdr[1].text = "Revenue"
    hdr[2].text = "Expenses"
    table.rows[1].cells[0].text = "Q1"
    table.rows[1].cells[1].text = "$142,000"
    table.rows[1].cells[2].text = "$118,000"
    table.rows[2].cells[0].text = "Q2"
    table.rows[2].cells[1].text = "$158,000"
    table.rows[2].cells[2].text = "$131,000"
    doc.add_heading("3. Technology", level=2)
    doc.add_paragraph(
        "Website accessibility audit completed. "
        "See https://bits-acb.org for details.")
    doc.save(str(path))
    return path


def _make_xlsx(path: Path) -> bool:
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Budget 2026"
        for cell, val in zip(["A1", "B1", "C1", "D1"],
                              ["Quarter", "Revenue", "Expenses", "Net"]):
            ws[cell] = val
        for row in [("Q1 2026", 142000, 118000, 24000),
                    ("Q2 2026", 158000, 131000, 27000),
                    ("Q3 2026", 155000, 127500, 27500),
                    ("Q4 2026", 161000, 134000, 27000)]:
            ws.append(row)
        ws2 = wb.create_sheet("Members")
        ws2["A1"] = "Region"
        ws2["B1"] = "Count"
        for row in [("Northeast", 412), ("Southeast", 318),
                    ("Midwest", 274), ("West", 200)]:
            ws2.append(row)
        wb.save(str(path))
        return True
    except ImportError:
        return False


def _make_pptx(path: Path) -> bool:
    try:
        from pptx import Presentation
        prs = Presentation()
        layout = prs.slide_layouts[1]
        s1 = prs.slides.add_slide(layout)
        s1.shapes.title.text = "BITS Annual Report 2026"
        s1.placeholders[1].text = (
            "Membership: 1,204 active members\n"
            "Revenue: $616,000 annual\n"
            "Website: WCAG 2.2 AA compliant")
        s2 = prs.slides.add_slide(layout)
        s2.shapes.title.text = "Technology Committee"
        s2.placeholders[1].text = (
            "Website audit: 14 missing alt text.\n"
            "3 color-only violations.\n"
            "2 form fields without labels.")
        s3 = prs.slides.add_slide(layout)
        s3.shapes.title.text = "Financial Summary"
        s3.placeholders[1].text = (
            "Q1 Revenue: $142,000\n"
            "Q2 Revenue: $158,000\n"
            "Q3 Revenue: $155,000\n"
            "Q4 Revenue: $161,000")
        prs.save(str(path))
        return True
    except ImportError:
        return False


def _make_text_pdf(path: Path) -> bool:
    try:
        from weasyprint import HTML as WP_HTML
        html = textwrap.dedent("""\
            <!DOCTYPE html>
            <html lang="en">
            <head><meta charset="utf-8"><title>BITS Report 2026</title></head>
            <body>
            <h1>BITS Board Report 2026</h1>
            <h2>Financial Summary</h2>
            <p>The organization generated $616,000 in revenue for fiscal year 2026.</p>
            <table>
              <thead><tr><th>Quarter</th><th>Revenue</th><th>Expenses</th></tr></thead>
              <tbody>
                <tr><td>Q1</td><td>$142,000</td><td>$118,000</td></tr>
                <tr><td>Q2</td><td>$158,000</td><td>$131,000</td></tr>
              </tbody>
            </table>
            <h2>Technology</h2>
            <p>Website audit: 14 alt text issues, 3 color-only violations.</p>
            <p>See https://bits-acb.org for full details.</p>
            </body></html>
        """)
        WP_HTML(string=html).write_pdf(str(path))
        return True
    except Exception:
        return False


def _make_epub(path: Path) -> Path:
    import zipfile
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(zipfile.ZipInfo("mimetype"), "application/epub+zip")
        zf.writestr("META-INF/container.xml", textwrap.dedent("""\
            <?xml version="1.0"?>
            <container version="1.0"
                xmlns="urn:oasis:schemas:container">
              <rootfiles>
                <rootfile full-path="EPUB/content.opf"
                    media-type="application/oebps-package+xml"/>
              </rootfiles>
            </container>"""))
        zf.writestr("EPUB/content.opf", textwrap.dedent("""\
            <?xml version="1.0" encoding="utf-8"?>
            <package xmlns="http://www.idpf.org/2007/opf" version="3.0"
                unique-identifier="uid">
              <metadata xmlns:dc="http://purl.org/dc/elements/1.1/">
                <dc:identifier id="uid">bits-report-2026</dc:identifier>
                <dc:title>BITS Report 2026</dc:title>
                <dc:language>en</dc:language>
              </metadata>
              <manifest>
                <item id="chapter1" href="chapter1.xhtml"
                    media-type="application/xhtml+xml"/>
              </manifest>
              <spine><itemref idref="chapter1"/></spine>
            </package>"""))
        zf.writestr("EPUB/chapter1.xhtml", textwrap.dedent("""\
            <?xml version="1.0" encoding="utf-8"?>
            <!DOCTYPE html>
            <html xmlns="http://www.w3.org/1999/xhtml" lang="en">
            <head><title>BITS Report 2026</title></head>
            <body>
              <h1>BITS Board Report 2026</h1>
              <h2>Financial Summary</h2>
              <p>Annual revenue reached $616,000 in 2026. The organization is stable.</p>
              <h2>Technology Committee</h2>
              <p>The website audit found 14 images missing alt text.</p>
              <h2>Membership</h2>
              <p>Total active membership: 1,204 members across all chapters.</p>
            </body>
            </html>"""))
    return path


def _make_html_file(path: Path) -> Path:
    path.write_text(textwrap.dedent("""\
        <!DOCTYPE html>
        <html lang="en">
        <head><meta charset="utf-8"><title>BITS Newsletter Spring 2026</title></head>
        <body>
          <h1>BITS Newsletter -- Spring 2026</h1>
          <h2>President's Message</h2>
          <p>Welcome to the spring 2026 edition of the BITS newsletter.</p>
          <p>Membership has grown to 1,204 members this quarter.</p>
          <h2>Technology Updates</h2>
          <p>The GLOW toolkit now supports Word, Excel, PowerPoint, PDF, and ePub.</p>
          <h2>Upcoming Events</h2>
          <ul>
            <li>Annual Conference: June 14-16, 2026</li>
            <li>Regional Meetups: July through September</li>
          </ul>
          <p>See <a href="https://bits-acb.org/events">BITS Events Calendar</a>.</p>
        </body>
        </html>"""), encoding="utf-8")
    return path


def _make_csv_file(path: Path) -> Path:
    path.write_text(
        "Quarter,Revenue,Expenses,Net\n"
        "Q1 2026,$142000,$118000,$24000\n"
        "Q2 2026,$158000,$131000,$27000\n"
        "Q3 2026,$155000,$127500,$27500\n"
        "Q4 2026,$161000,$134000,$27000\n"
        "Total,$616000,$510500,$105500\n",
        encoding="utf-8")
    return path


def _markitdown_extract(file_path: Path) -> str:
    from markitdown import MarkItDown
    md = MarkItDown()
    result = md.convert(str(file_path))
    return result.text_content or ""


# ===========================================================================
# LAYER 0 -- Document ingestion
# ===========================================================================

def layer0(tmpdir: Path) -> None:
    print("\n" + "=" * 65)
    print("LAYER 0: Document ingestion")
    print("=" * 65)

    # 00a: .md direct read
    md_path = tmpdir / "test.md"
    md_path.write_text(RICH_BAD_DOC, encoding="utf-8")
    text = md_path.read_text(encoding="utf-8")
    ok = "BITS Board Meeting" in text and "Financial" in text
    _record("00a", ".md direct read_text()", ok, text[:300],
            {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")

    # 00b: .docx
    docx_path = tmpdir / "test.docx"
    try:
        _make_docx(docx_path)
        text = _markitdown_extract(docx_path)
        ok = len(text) > 50 and any(
            kw in text.lower() for kw in ["bits", "budget", "financial", "quarter"])
        _record("00b", ".docx MarkItDown extraction", ok, text[:300],
                {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")
    except Exception as exc:
        _record("00b", ".docx MarkItDown extraction", False, "",
                {"error": str(exc)}, layer=0, notes=str(exc))

    # 00c: .xlsx
    xlsx_path = tmpdir / "test.xlsx"
    if _make_xlsx(xlsx_path):
        try:
            text = _markitdown_extract(xlsx_path)
            ok = len(text) > 10 and any(
                kw in text.lower() for kw in ["quarter", "revenue", "q1"])
            _record("00c", ".xlsx MarkItDown extraction", ok, text[:300],
                    {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")
        except Exception as exc:
            _record("00c", ".xlsx MarkItDown extraction", False, "",
                    {"error": str(exc)}, layer=0, notes=str(exc))
    else:
        print("  [SKIP] 00c: .xlsx (openpyxl not installed)")

    # 00d: .pptx
    pptx_path = tmpdir / "test.pptx"
    if _make_pptx(pptx_path):
        try:
            text = _markitdown_extract(pptx_path)
            ok = len(text) > 20 and any(
                kw in text.lower() for kw in ["bits", "technology", "membership", "annual"])
            _record("00d", ".pptx MarkItDown extraction", ok, text[:300],
                    {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")
        except Exception as exc:
            _record("00d", ".pptx MarkItDown extraction", False, "",
                    {"error": str(exc)}, layer=0, notes=str(exc))
    else:
        print("  [SKIP] 00d: .pptx (python-pptx not installed)")

    # 00e: .pdf (text-based)
    pdf_path = tmpdir / "test.pdf"
    if _make_text_pdf(pdf_path):
        try:
            text = _markitdown_extract(pdf_path)
            ok = len(text) > 20 and any(
                kw in text.lower() for kw in ["bits", "financial", "revenue"])
            _record("00e", ".pdf text-based MarkItDown", ok, text[:300],
                    {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")
        except Exception as exc:
            _record("00e", ".pdf text-based MarkItDown", False, "",
                    {"error": str(exc)}, layer=0, notes=str(exc))
    else:
        print("  [SKIP] 00e: .pdf (WeasyPrint unavailable)")

    # 00f: .epub
    epub_path = tmpdir / "test.epub"
    _make_epub(epub_path)
    try:
        text = _markitdown_extract(epub_path)
        ok = len(text) > 20 and any(
            kw in text.lower() for kw in ["bits", "financial", "membership"])
        _record("00f", ".epub MarkItDown extraction", ok, text[:300],
                {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")
    except Exception as exc:
        _record("00f", ".epub MarkItDown extraction", False, "",
                {"error": str(exc)}, layer=0, notes=str(exc))

    # 00g: .html
    html_path = tmpdir / "test.html"
    _make_html_file(html_path)
    try:
        text = _markitdown_extract(html_path)
        ok = len(text) > 20 and "BITS" in text
        _record("00g", ".html MarkItDown extraction", ok, text[:300],
                {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")
    except Exception as exc:
        _record("00g", ".html MarkItDown extraction", False, "",
                {"error": str(exc)}, layer=0, notes=str(exc))

    # 00h: .csv
    csv_path = tmpdir / "test.csv"
    _make_csv_file(csv_path)
    try:
        text = _markitdown_extract(csv_path)
        ok = len(text) > 10 and any(kw in text.lower() for kw in ["quarter", "q1", "revenue"])
        _record("00h", ".csv MarkItDown extraction", ok, text[:300],
                {"chars": len(text)}, layer=0, notes=f"{len(text)} chars")
    except Exception as exc:
        _record("00h", ".csv MarkItDown extraction", False, "",
                {"error": str(exc)}, layer=0, notes=str(exc))

    # 00i: long-doc truncation
    result = _build_document_context("What are the financial findings?", LONG_DOC)
    ok = len(result) <= 8000
    _record("00i", "_build_document_context long-doc truncation", ok,
            f"In:{len(LONG_DOC)} Out:{len(result)}\n\n{result[:200]}",
            {"input": len(LONG_DOC), "output": len(result)}, layer=0,
            notes=f"{len(LONG_DOC)} -> {len(result)} chars")

    # 00j: short-doc passthrough
    short = PLAIN_LANGUAGE_DOC
    result = _build_document_context("What did we decide?", short)
    ok = result.strip() == short.strip()
    _record("00j", "_build_document_context short-doc passthrough", ok,
            f"In:{len(short)} Out:{len(result)}",
            {"match": ok}, layer=0)

    # 00k: keyword snippet extraction
    result = _build_document_context("What is the budget and revenue?", LONG_DOC)
    ok = ("budget" in result.lower() or "revenue" in result.lower()
          or "financial" in result.lower())
    _record("00k", "_build_document_context keyword snippets", ok,
            f"Out:{len(result)}\n\n{result[:300]}",
            {"has_kw": ok}, layer=0,
            notes=f"keyword-relevant snippets in {len(result)} chars")


# ===========================================================================
# LAYER 1 -- Raw OpenRouter gateway
# ===========================================================================

def layer1() -> None:
    print("\n" + "=" * 65)
    print("LAYER 1: Raw OpenRouter gateway")
    print("=" * 65)

    # 01: primary model (gpt-4o-mini)
    try:
        data = _chat_raw(
            "openai/gpt-4o-mini",
            [
                {"role": "system", "content": "Answer in 1-2 sentences."},
                {"role": "user",
                 "content": "What font does the ACB Large Print Guidelines require for body text?"},
            ],
        )
        answer = data["choices"][0]["message"]["content"].strip()
        ok = bool(answer) and any(kw in answer.lower() for kw in ["arial", "sans", "font"])
        _record("01", "primary-model-ACB-font-question", ok, answer, data, layer=1)
    except Exception as exc:
        _record("01", "primary-model-ACB-font-question", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))

    # 02: paid model echo
    try:
        data = _chat_raw(
            "openai/gpt-4o-mini",
            [{"role": "user", "content": "Reply with exactly: GLOW_PAID_OK"}],
            max_tokens=20,
        )
        answer = data["choices"][0]["message"]["content"].strip()
        _record("02", "paid-model-echo", "GLOW_PAID_OK" in answer, answer, data, layer=1)
    except Exception as exc:
        _record("02", "paid-model-echo", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))

    # 03: gateway chat() free-first
    try:
        answer, esc = gateway_chat(
            question="What is WCAG 2.2 AA in one sentence?",
            system_prompt="You are the GLOW accessibility assistant. Be concise.",
            session_hash=SESSION,
        )
        _record("03", "gateway-chat-free-first", bool(answer), answer,
                {"escalated": esc}, layer=1, notes=f"escalated={esc}")
    except Exception as exc:
        _record("03", "gateway-chat-free-first", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))

    # 04: gateway chat() with document context
    try:
        answer, esc = gateway_chat(
            question="What accessibility issues does this document have?",
            system_prompt="You are the GLOW auditor. List specific issues.",
            session_hash=SESSION,
            document_text=RICH_BAD_DOC[:3000],
        )
        ok = bool(answer) and any(kw in answer.lower() for kw in
                                   ["italic", "bold", "alt", "font", "issue", "heading", "link"])
        _record("04", "gateway-chat-document-context", ok, answer,
                {"escalated": esc}, layer=1, notes=f"escalated={esc}")
    except Exception as exc:
        _record("04", "gateway-chat-document-context", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))

    # 05: multi-turn history
    history = [
        {"role": "user", "content": "My document uses 10pt Times New Roman."},
        {"role": "assistant", "content": "ACB requires 18pt Arial for body text."},
        {"role": "user", "content": "What about headings?"},
        {"role": "assistant", "content": "Headings must be 22pt Arial."},
    ]
    try:
        answer, esc = gateway_chat(
            question="And subheadings -- what size?",
            system_prompt="You are the GLOW accessibility assistant. Be concise.",
            session_hash=SESSION,
            conversation_history=history,
        )
        ok = bool(answer) and any(kw in answer.lower() for kw in
                                   ["20", "pt", "point", "subhead"])
        _record("05", "gateway-chat-multi-turn", ok, answer,
                {"escalated": esc, "turns": len(history)}, layer=1,
                notes=f"escalated={esc}")
    except Exception as exc:
        _record("05", "gateway-chat-multi-turn", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))

    # 06: rich ACB audit
    try:
        answer, esc = gateway_chat(
            question=(
                "Audit this HTML for ACB Large Print compliance and list each violation "
                "with rule ID and severity:\n\n"
                "<p><em>Minutes from April 8</em></p>\n"
                "<p style=\"text-align:center\"><strong>All items approved</strong></p>\n"
                "<p><a href=\"https://internal.example.com\">click here</a></p>"
            ),
            system_prompt=(
                "You are the GLOW ACB auditor. Format each finding as: "
                "[RULE-ID] Description (Severity: critical/high/medium/low)."
            ),
            session_hash=SESSION,
        )
        ok = bool(answer) and any(kw in answer.lower() for kw in
                                   ["italic", "acb", "click", "center", "align"])
        _record("06", "gateway-chat-rich-ACB-audit", ok, answer,
                {"escalated": esc}, layer=1, notes=f"escalated={esc}")
    except Exception as exc:
        _record("06", "gateway-chat-rich-ACB-audit", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))

    # 07: uncertainty trigger
    try:
        answer, esc = gateway_chat(
            question="What pixel resolution does the 2031 ACB revision require?",
            system_prompt="Answer only from known facts. Say 'I don't know' if uncertain.",
            session_hash=SESSION,
        )
        _record("07", "gateway-chat-uncertainty", bool(answer), answer,
                {"escalated": esc}, layer=1,
                notes=f"escalated={esc} (acceptable)")
    except Exception as exc:
        _record("07", "gateway-chat-uncertainty", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))

    # 08: quota status
    try:
        status = get_quota_status(SESSION)
        ok = all(k in status for k in
                 ("ai_configured", "budget_usd", "chat_available", "audio_available"))
        _record("08", "quota-status-no-model-call", ok,
                json.dumps(status, indent=2), status, layer=1,
                notes=f"remaining=${status.get('budget_remaining_usd', '?')}")
    except Exception as exc:
        _record("08", "quota-status-no-model-call", False, "",
                {"error": str(exc)}, layer=1, notes=str(exc))


# ===========================================================================
# LAYER 2 -- Document tools
# ===========================================================================

def layer2() -> None:
    print("\n" + "=" * 65)
    print("LAYER 2: Document tools (ToolRegistry)")
    print("=" * 65)

    # 09: get_document_stats
    _, t = _ctx(RICH_BAD_DOC, "board-minutes.md")
    r = t.get_document_stats()
    ok = any(kw in r for kw in ["Words:", "Headings:", "Tables:", "Reading time:"])
    _record("09", "tool-get-document-stats", ok, r, {"raw": r}, layer=2)

    # 10: list_headings
    _, t = _ctx(RICH_BAD_DOC, "board-minutes.md")
    r = t.list_headings()
    ok = "BITS Board Meeting" in r and "Financial" in r and "Technology" in r
    _record("10", "tool-list-headings-rich", ok, r, {"raw": r}, layer=2)

    # 11: check_heading_hierarchy valid
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.check_heading_hierarchy()
    ok = "valid" in r.lower() or "no skipped" in r.lower()
    _record("11", "tool-heading-hierarchy-valid", ok, r, {"raw": r}, layer=2)

    # 12: check_heading_hierarchy H2->H4 skip
    _, t = _ctx(SKIPPED_HEADING_DOC)
    r = t.check_heading_hierarchy()
    ok = "skipped" in r.lower() or "issue" in r.lower() or "h2" in r.lower()
    _record("12", "tool-heading-hierarchy-skipped", ok, r, {"raw": r}, layer=2)

    # 13: find_faux_headings violations
    faux = ("# Real Heading\n\n**This Is A Faux Heading**\n\nBody text.\n\n"
            "**Another Fake Section Title**\n\nMore content.\n")
    _, t = _ctx(faux)
    r = t.find_faux_headings()
    ok = any(kw in r.lower() for kw in ["faux", "fake", "bold", "Faux"])
    _record("13", "tool-faux-headings-violations", ok, r, {"raw": r}, layer=2)

    # 14: find_faux_headings clean
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.find_faux_headings()
    ok = "no obvious" in r.lower() or "not found" in r.lower() or "0" in r
    _record("14", "tool-faux-headings-clean", ok, r, {"raw": r}, layer=2)

    # 15: check_list_structure deep nesting
    deep = "# Doc\n\n- Item\n  - Nested\n    - Deep\n      - Too deep\n        - Way too deep\n"
    _, t = _ctx(deep)
    r = t.check_list_structure()
    ok = any(kw in r.lower() for kw in ["nesting", "deep", "level"])
    _record("15", "tool-list-structure-deep", ok, r, {"raw": r}, layer=2)

    # 16: check_list_structure clean
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.check_list_structure()
    ok = any(kw in r.lower() for kw in ["bullet", "numbered", "item", "list"])
    _record("16", "tool-list-structure-clean", ok, r, {"raw": r}, layer=2)

    # 17: estimate_reading_order
    _, t = _ctx(RICH_BAD_DOC)
    r = t.estimate_reading_order()
    ok = any(kw in r.lower() for kw in ["table", "risk", "reading"])
    _record("17", "tool-estimate-reading-order", ok, r, {"raw": r}, layer=2)

    # 18: find_section match
    _, t = _ctx(RICH_BAD_DOC)
    r = t.find_section("Financial")
    ok = "financial" in r.lower() and len(r) > 20
    _record("18", "tool-find-section-match", ok, r[:300], {"raw": r[:300]}, layer=2)

    # 19: find_section no match
    _, t = _ctx(RICH_BAD_DOC)
    r = t.find_section("Nonexistent Section XYZ")
    ok = "not found" in r.lower()
    _record("19", "tool-find-section-no-match", ok, r, {"raw": r}, layer=2)

    # 20: search_text match -- search for "reserve" which appears literally in RICH_BAD_DOC
    _, t = _ctx(RICH_BAD_DOC)
    r = t.search_text("reserve")
    ok = "reserve" in r.lower() and "line" in r.lower()
    _record("20", "tool-search-text-match", ok, r, {"raw": r}, layer=2)

    # 21: search_text no match
    _, t = _ctx(RICH_BAD_DOC)
    r = t.search_text("xyzzy_not_present")
    ok = "no matches" in r.lower() or "not found" in r.lower()
    _record("21", "tool-search-text-no-match", ok, r, {"raw": r}, layer=2)

    # 22: extract_table index 0
    _, t = _ctx(RICH_BAD_DOC)
    r = t.extract_table("0")
    ok = any(kw in r.lower() for kw in ["quarter", "revenue", "expenses", "|"])
    _record("22", "tool-extract-table-index-0", ok, r, {"raw": r}, layer=2)

    # 23: extract_table not found
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.extract_table("99")
    ok = "not found" in r.lower() or "available" in r.lower()
    _record("23", "tool-extract-table-not-found", ok, r, {"raw": r}, layer=2)

    # 24: get_images
    _, t = _ctx(RICH_BAD_DOC)
    r = t.get_images()
    ok = bool(r)
    _record("24", "tool-get-images", ok, r, {"raw": r}, layer=2)

    # 24b: get_document_summary
    _, t = _ctx(RICH_BAD_DOC)
    r = t.get_document_summary()
    ok = bool(r) and any(kw in r.lower() for kw in ["document", "type", "section", "title", "meeting", "words"])
    _record("24b", "tool-get-document-summary", ok, r, {"raw": r}, layer=2)

    # 24c: get_section_content (existing section)
    _, t = _ctx(RICH_BAD_DOC)
    r = t.get_section_content("financial")
    ok = bool(r) and "not found" not in r.lower()
    _record("24c", "tool-get-section-content-found", ok, r, {"raw": r}, layer=2)

    # 24d: get_section_content (missing section)
    _, t = _ctx(RICH_BAD_DOC)
    r = t.get_section_content("nonexistent-xyz-section")
    ok = "not found" in r.lower() or "available" in r.lower()
    _record("24d", "tool-get-section-content-not-found", ok, r, {"raw": r}, layer=2)

    # 24e: get_decisions_and_actions
    _, t = _ctx(RICH_BAD_DOC)
    r = t.get_decisions_and_actions()
    ok = bool(r)
    _record("24e", "tool-get-decisions-and-actions", ok, r, {"raw": r}, layer=2)

    # 24f: get_what_passes
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.get_what_passes()
    ok = bool(r) and "pass" in r.lower()
    _record("24f", "tool-get-what-passes", ok, r, {"raw": r}, layer=2)


# ===========================================================================
# LAYER 3 -- Compliance Agent tools
# ===========================================================================

def layer3() -> None:
    print("\n" + "=" * 65)
    print("LAYER 3: Compliance Agent tools")
    print("=" * 65)

    # 25: audit violations
    _, t = _ctx(RICH_BAD_DOC)
    r = t.run_accessibility_audit()
    ok = bool(r) and any(kw in r.lower() for kw in
                          ["italic", "bold", "url", "issue", "audit", "compliance"])
    _record("25", "compliance-audit-violations", ok, r, {"raw": r}, layer=3)

    # 26: audit clean
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.run_accessibility_audit()
    ok = bool(r)
    _record("26", "compliance-audit-clean", ok, r, {"raw": r}, layer=3)

    # 27: get_compliance_score
    _, t = _ctx(RICH_BAD_DOC)
    t.run_accessibility_audit()
    r = t.get_compliance_score()
    ok = bool(r)
    _record("27", "compliance-get-score", ok, r, {"raw": r}, layer=3)

    # 28: get_critical_findings
    _, t = _ctx(RICH_BAD_DOC)
    t.run_accessibility_audit()
    r = t.get_critical_findings()
    ok = bool(r)
    _record("28", "compliance-get-critical", ok, r, {"raw": r}, layer=3)

    # 29: get_auto_fixable_findings
    _, t = _ctx(RICH_BAD_DOC)
    t.run_accessibility_audit()
    r = t.get_auto_fixable_findings()
    ok = bool(r)
    _record("29", "compliance-get-auto-fixable", ok, r, {"raw": r}, layer=3)


# ===========================================================================
# LAYER 4 -- Content Agent tools
# ===========================================================================

def layer4() -> None:
    print("\n" + "=" * 65)
    print("LAYER 4: Content Agent tools")
    print("=" * 65)

    # 30: emphasis violations
    _, t = _ctx(RICH_BAD_DOC)
    r = t.check_emphasis_patterns()
    ok = any(kw in r.lower() for kw in ["italic", "critical", "bold", "acb-no-italic"])
    _record("30", "content-emphasis-violations", ok, r, {"raw": r}, layer=4)

    # 31: emphasis clean
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.check_emphasis_patterns()
    ok = "ok" in r.lower() or "no italic" in r.lower()
    _record("31", "content-emphasis-clean", ok, r, {"raw": r}, layer=4)

    # 32: link text violations
    _, t = _ctx(RICH_BAD_DOC)
    r = t.check_link_text()
    ok = any(kw in r.lower() for kw in ["bare", "url", "generic", "click here", "link"])
    _record("32", "content-link-text-violations", ok, r, {"raw": r}, layer=4)

    # 33: link text clean
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.check_link_text()
    ok = "ok" in r.lower() or "no bare" in r.lower()
    _record("33", "content-link-text-clean", ok, r, {"raw": r}, layer=4)

    # 34: reading level complex
    _, t = _ctx(RICH_BAD_DOC)
    r = t.check_reading_level()
    ok = any(kw in r.lower() for kw in ["sentence", "word", "level", "reading"])
    _record("34", "content-reading-level-complex", ok, r, {"raw": r}, layer=4)

    # 35: reading level plain
    _, t = _ctx(PLAIN_LANGUAGE_DOC)
    r = t.check_reading_level()
    ok = "level" in r.lower() and any(kw in r.lower() for kw in ["low", "plain", "sentence"])
    _record("35", "content-reading-level-plain", ok, r, {"raw": r}, layer=4)

    # 36: alignment center found
    _, t = _ctx(RICH_BAD_DOC)
    r = t.check_alignment_hints()
    ok = "center" in r.lower() or "alignment" in r.lower()
    _record("36", "content-alignment-center", ok, r, {"raw": r}, layer=4)

    # 37: alignment clean
    _, t = _ctx(RICH_CLEAN_DOC)
    r = t.check_alignment_hints()
    ok = "no explicit" in r.lower() or "ok" in r.lower() or "not detected" in r.lower()
    _record("37", "content-alignment-clean", ok, r, {"raw": r}, layer=4)


# ===========================================================================
# LAYER 5 -- Remediation Agent tools
# ===========================================================================

def layer5() -> None:
    print("\n" + "=" * 65)
    print("LAYER 5: Remediation Agent tools")
    print("=" * 65)

    _, t = _ctx(RICH_BAD_DOC)

    # 38-41: explain_rule known rules
    for pid, rule_id in [("38", "ACB-NO-ITALIC"), ("39", "ACB-FAUX-HEADING"),
                          ("40", "ACB-ALIGNMENT"), ("41", "ACB-TABLE-HEADERS")]:
        r = t.explain_rule(rule_id)
        ok = "fix" in r.lower() and rule_id in r
        _record(pid, f"remediation-explain-{rule_id}", ok, r, {"raw": r}, layer=5)

    # 42: explain_rule unknown
    r = t.explain_rule("ACB-DOES-NOT-EXIST")
    ok = "not found" in r.lower() or "known rules" in r.lower()
    _record("42", "remediation-explain-unknown-rule", ok, r, {"raw": r}, layer=5)

    # 43: suggest_fix
    r = t.suggest_fix("ACB-FONT-SIZE-BODY")
    ok = "fix" in r.lower() and ("18" in r or "font" in r.lower())
    _record("43", "remediation-suggest-fix-font", ok, r, {"raw": r}, layer=5)

    # 44: prioritize_findings
    _, t2 = _ctx(RICH_BAD_DOC)
    r = t2.prioritize_findings()
    ok = bool(r) and any(kw in r.lower() for kw in
                          ["critical", "fix", "1.", "priority", "italic"])
    _record("44", "remediation-prioritize-heuristic", ok, r, {"raw": r}, layer=5)

    # 45: estimate_fix_impact
    _, t3 = _ctx(RICH_BAD_DOC)
    r = t3.estimate_fix_impact()
    ok = bool(r)
    _record("45", "remediation-estimate-impact-no-cache", ok, r, {"raw": r}, layer=5)

    # 46: alt text 2 of 3 missing
    partial = "# Doc\n\n![Good desc](a.jpg)\n![](b.jpg)\n![](c.jpg)\n"
    _, t4 = _ctx(partial)
    r = t4.check_image_alt_text()
    ok = "missing" in r.lower() or "empty" in r.lower()
    _record("46", "remediation-alt-text-partial-missing", ok, r, {"raw": r}, layer=5)

    # 47: alt text all present
    _, t5 = _ctx(RICH_CLEAN_DOC)
    r = t5.check_image_alt_text()
    ok = "ok" in r.lower() or "all" in r.lower()
    _record("47", "remediation-alt-text-all-present", ok, r, {"raw": r}, layer=5)

    # 48: alt text no images
    _, t6 = _ctx(PLAIN_LANGUAGE_DOC)
    r = t6.check_image_alt_text()
    ok = "no markdown images" in r.lower() or "no" in r.lower()
    _record("48", "remediation-alt-text-no-images", ok, r, {"raw": r}, layer=5)


# ===========================================================================
# LAYER 6 -- Pre-flight keyword dispatcher
# ===========================================================================

def layer6() -> None:
    print("\n" + "=" * 65)
    print("LAYER 6: Pre-flight keyword dispatcher")
    print("=" * 65)

    def dispatch(q: str, doc: str = RICH_BAD_DOC) -> str:
        _, t = _ctx(doc)
        return t.dispatch_for_question(q)

    cases = [
        ("49", "dispatch-audit-compliance",
         "What violations and compliance findings does this document have?",
         ["audit", "compliance", "findings", "score"]),
        ("50", "dispatch-fix-remediation",
         "How do I fix all the italic and bold issues in this document?",
         ["fix", "prioritize", "remediat", "impact"]),
        ("51", "dispatch-heading-structure",
         "What is the heading structure and hierarchy of this document?",
         ["heading", "h1", "h2", "hierarchy", "outline"]),
        ("52", "dispatch-emphasis-italic",
         "Are there any italic or bold emphasis formatting violations?",
         ["italic", "emphasis", "bold", "acb-no-italic"]),
        ("53", "dispatch-link-url",
         "Check all link text and bare URLs in this document.",
         ["link", "url", "bare", "click"]),
        ("54", "dispatch-image-alt",
         "Which images are missing alt text in this document?",
         ["alt", "image", "missing", "img"]),
        ("55", "dispatch-list-bullet",
         "Check all bullet lists and numbered lists for deep nesting.",
         ["list", "bullet", "item", "nesting"]),
        ("56", "dispatch-reading-level",
         "What is the reading level and sentence complexity of this document?",
         ["sentence", "reading", "level", "word"]),
        ("57", "dispatch-find-section",
         "Show me the section on Financial Summary from this document.",
         ["financial", "section", "summary"]),
        ("58", "dispatch-search-term",
         "Search for all places where italic appears in this document.",
         ["italic", "line", "search", "match"]),
        ("59", "dispatch-table-data",
         "How many tables and data columns does this document contain?",
         ["table", "reading order", "risk"]),
        ("60", "dispatch-acb-rule-id",
         "Explain ACB-NO-ITALIC and the steps to fix it in my document.",
         ["acb-no-italic", "italic", "fix"]),
        ("61", "dispatch-no-keywords-default",
         "Hello, what can you tell me about this document overall?",
         ["words", "document", "heading", "stats"]),
        ("62", "dispatch-combined-italic-heading",
         "Are there italic violations and heading structure issues?",
         ["italic", "emphasis", "heading", "hierarchy"]),
    ]

    for path_id, name, question, expected_kws in cases:
        r = dispatch(question)
        ok = any(kw in r.lower() for kw in expected_kws)
        tools_fired = [
            ln.strip("[]")
            for ln in r.splitlines()
            if ln.startswith("[") and "===" not in ln
        ]
        _record(path_id, name, ok, r[:600],
                {"tools_fired": tools_fired, "question": question},
                layer=6, notes=f"tools={tools_fired}")


# ===========================================================================
# LAYER 7 -- Full GLOW pipeline
# ===========================================================================

def layer7() -> None:
    print("\n" + "=" * 65)
    print("LAYER 7: Full GLOW pipeline")
    print("=" * 65)

    questions = [
        ("63", "pipeline-all-critical-violations",
         "What are all the critical accessibility violations in this document? "
         "List each one with its ACB rule ID.",
         ["italic", "acb", "violation", "critical", "heading", "link", "bold"]),
        ("64", "pipeline-prioritize-fixes",
         "Prioritize all the accessibility fixes from most to least impactful "
         "and tell me which ones can be auto-fixed.",
         ["fix", "critical", "auto", "priority", "italic", "remediat"]),
        ("65", "pipeline-heading-structure",
         "What is the complete heading structure of this document? "
         "Are any heading levels skipped?",
         ["heading", "h1", "h2", "h3", "structure", "hierarchy"]),
        ("66", "pipeline-explain-italic-fix",
         "Explain the italic violation in detail and give step-by-step "
         "instructions to fix it.",
         ["italic", "fix", "acb-no-italic", "underline", "emphasis"]),
        ("67", "pipeline-table-check",
         "What tables are in this document and do they have proper header rows?",
         ["table", "header", "quarter", "revenue", "financial"]),
        ("68", "pipeline-reading-level",
         "What is the reading level of this document? Is the language plain "
         "enough for people with cognitive disabilities?",
         ["reading", "sentence", "level", "word", "complex", "plain"]),
        ("69", "pipeline-section-budget",
         "What section discusses the budget and what does it say about reserves?",
         ["budget", "financial", "reserve", "section", "quarter"]),
    ]

    for path_id, name, question, expected_kws in questions:
        try:
            answer, esc, tools_run = _full_pipeline(question)
            ok = (bool(answer) and len(answer) > 40 and
                  any(kw in answer.lower() for kw in expected_kws))
            _record(path_id, name, ok, answer,
                    {"escalated": esc, "tools": tools_run},
                    layer=7, notes=f"tools={tools_run} escalated={esc}")
        except Exception as exc:
            _record(path_id, name, False, "", {"error": str(exc)},
                    layer=7, notes=str(exc))

    # 70: multi-turn follow-up
    history = [
        {"role": "user",
         "content": "What italic issues are in this document?"},
        {"role": "assistant",
         "content": (
             "The document has two italic violations: "
             "'Minutes recorded by the Secretary' and "
             "'The committee chair reviewed each report in order.' "
             "Both violate ACB-NO-ITALIC (critical severity)."
         )},
    ]
    try:
        _, tools = _ctx(RICH_BAD_DOC)
        preflight = tools.dispatch_for_question(
            "How do I fix the italic violations you mentioned?")
        full_system = f"{_CHAT_SYSTEM_PROMPT}\n\n{preflight}"
        answer, esc = gateway_chat(
            question="How do I fix the italic violations you mentioned?",
            system_prompt=full_system,
            session_hash=SESSION,
            document_text=_build_document_context("italic fix", RICH_BAD_DOC),
            conversation_history=history,
        )
        ok = bool(answer) and any(kw in answer.lower() for kw in
                                   ["italic", "fix", "underline", "remove", "acb"])
        _record("70", "pipeline-multi-turn-follow-up", ok, answer,
                {"escalated": esc, "turns": len(history)},
                layer=7, notes=f"escalated={esc}")
    except Exception as exc:
        _record("70", "pipeline-multi-turn-follow-up", False, "",
                {"error": str(exc)}, layer=7, notes=str(exc))

    # 70b: document summary tool is selected and used in pipeline
    try:
        answer, esc, tools_run = _full_pipeline("What is this document about and what topics does it cover?")
        ok = (bool(answer) and
              any(kw in answer.lower() for kw in ["meeting", "board", "document", "section", "report", "agenda", "bits"]))
        _record("70b", "pipeline-document-summary",
                ok, answer,
                {"escalated": esc, "tools": tools_run},
                layer=7, notes=f"tools={tools_run} escalated={esc}")
    except Exception as exc:
        _record("70b", "pipeline-document-summary", False, "",
                {"error": str(exc)}, layer=7, notes=str(exc))

    # 70c: section content tool used when user asks what a section says
    try:
        answer, esc, tools_run = _full_pipeline("What does the financial report say about reserves?")
        ok = (bool(answer) and
              any(kw in answer.lower() for kw in ["reserve", "budget", "financial", "354", "q1", "revenue"]))
        _record("70c", "pipeline-section-content",
                ok, answer,
                {"escalated": esc, "tools": tools_run},
                layer=7, notes=f"tools={tools_run} escalated={esc}")
    except Exception as exc:
        _record("70c", "pipeline-section-content", False, "",
                {"error": str(exc)}, layer=7, notes=str(exc))

    # 70d: decisions / action items extraction
    try:
        answer, esc, tools_run = _full_pipeline("What did the board vote on or decide?")
        ok = (bool(answer) and
              any(kw in answer.lower() for kw in ["voted", "approved", "unanimous", "decision", "motion", "resolved", "action"]))
        _record("70d", "pipeline-decisions-and-actions",
                ok, answer,
                {"escalated": esc, "tools": tools_run},
                layer=7, notes=f"tools={tools_run} escalated={esc}")
    except Exception as exc:
        _record("70d", "pipeline-decisions-and-actions", False, "",
                {"error": str(exc)}, layer=7, notes=str(exc))


# ===========================================================================
# LAYER 8 -- Vision and audio gateway
# ===========================================================================

def layer8(tmpdir: Path) -> None:
    print("\n" + "=" * 65)
    print("LAYER 8: Vision and audio gateway")
    print("=" * 65)

    # 71: annotated PNG with text
    png_annotated = _make_png_with_text("ACB Large Print: Arial 18pt minimum")
    try:
        result = describe_image(
            image_bytes=png_annotated,
            mime_type="image/png",
            prompt=(
                "Describe this image for accessibility review. "
                "Extract any visible text. Note whether this appears to be a "
                "document, chart, photo, or other content type."
            ),
            session_hash=SESSION,
        )
        ok = isinstance(result, str) and bool(result)
        _record("71", "vision-annotated-png-with-text", ok, result,
                {"chars": len(result)}, layer=8)
    except Exception as exc:
        _record("71", "vision-annotated-png-with-text", False, "",
                {"error": str(exc)}, layer=8, notes=str(exc))

    # 72: tiny synthetic PNG
    tiny = _make_tiny_png()
    try:
        result = describe_image(
            image_bytes=tiny,
            mime_type="image/png",
            prompt="Describe this image. Extract any visible text.",
            session_hash=SESSION,
        )
        ok = isinstance(result, str) and bool(result)
        _record("72", "vision-tiny-synthetic-png", ok, result,
                {"chars": len(result)}, layer=8)
    except Exception as exc:
        _record("72", "vision-tiny-synthetic-png", False, "",
                {"error": str(exc)}, layer=8, notes=str(exc))

    # 73: silent WAV
    wav_path = tmpdir / "silent.wav"
    _make_silent_wav(wav_path)
    try:
        result = transcribe(audio_path=wav_path, language="en", session_hash=SESSION)
        ok = isinstance(result, str)
        note = ("silent audio -> empty transcript is correct"
                if not result else f"got: {result[:80]!r}")
        _record("73", "audio-silent-wav", ok, result or "(empty)",
                {"chars": len(result)}, layer=8, notes=note)
    except RuntimeError as exc:
        _record("73", "audio-silent-wav", False, "",
                {"error": str(exc)}, layer=8, notes=f"RuntimeError: {exc}")
    except Exception as exc:
        _record("73", "audio-silent-wav", False, "",
                {"error": str(exc)}, layer=8, notes=str(exc))

    # 74: real mp3 (optional)
    mp3 = PROBE_ROOT / "reagan-30s.mp3"
    if not mp3.exists():
        print("  [SKIP] 74: audio-real-mp3 (reagan-30s.mp3 not found)")
        return
    try:
        result = transcribe(audio_path=mp3, language="en", session_hash=SESSION)
        ok = isinstance(result, str) and len(result) > 20
        _record("74", "audio-real-mp3", ok, result,
                {"chars": len(result)}, layer=8, notes=f"{len(result)} chars")
    except RuntimeError as exc:
        _record("74", "audio-real-mp3", False, "",
                {"error": str(exc)}, layer=8, notes=f"RuntimeError: {exc}")
    except Exception as exc:
        _record("74", "audio-real-mp3", False, "",
                {"error": str(exc)}, layer=8, notes=str(exc))


# ===========================================================================
# Main
# ===========================================================================

def main() -> None:
    parser = argparse.ArgumentParser(description="GLOW all-layers chat probe")
    parser.add_argument("--layer", type=int, default=None,
                        help="Run only this layer (0-8). Omit to run all layers.")
    args = parser.parse_args()

    if not API_KEY:
        print("ERROR: OPENROUTER_API_KEY not set.")
        sys.exit(1)

    run_all = args.layer is None
    only = args.layer

    print("=" * 65)
    print("GLOW All-Layers Chat Probe")
    print(f"Running: {'all 9 layers' if run_all else f'layer {only} only'}")
    print(f"Output:  {OUT_DIR}")
    print("=" * 65)

    with tempfile.TemporaryDirectory() as tmpdir_str:
        tmpdir = Path(tmpdir_str)

        if run_all or only == 0:
            layer0(tmpdir)
        if run_all or only == 1:
            layer1()
        if run_all or only == 2:
            layer2()
        if run_all or only == 3:
            layer3()
        if run_all or only == 4:
            layer4()
        if run_all or only == 5:
            layer5()
        if run_all or only == 6:
            layer6()
        if run_all or only == 7:
            layer7()
        if run_all or only == 8:
            layer8(tmpdir)

    summary_path = OUT_DIR / "summary.json"
    summary_path.write_text(json.dumps(RESULTS, indent=2, default=str), encoding="utf-8")

    layer_names = {
        0: "Document ingestion",
        1: "Raw gateway",
        2: "Document tools",
        3: "Compliance Agent",
        4: "Content Agent",
        5: "Remediation Agent",
        6: "Pre-flight dispatcher",
        7: "Full GLOW pipeline",
        8: "Vision + audio",
    }

    passed = sum(1 for r in RESULTS if r["ok"])
    total = len(RESULTS)

    print("\n" + "=" * 65)
    print(f"RESULTS: {passed}/{total} paths PASSED")
    print("-" * 65)
    for lnum in sorted(LAYER_COUNTERS):
        lc = LAYER_COUNTERS[lnum]
        lname = layer_names.get(lnum, f"Layer {lnum}")
        ltotal = lc["pass"] + lc["fail"]
        print(f"  Layer {lnum} ({lname}): {lc['pass']}/{ltotal}")
    print("-" * 65)
    print(f"Summary JSON: {summary_path}")
    print("=" * 65)

    if passed < total:
        print("\nFailed paths:")
        for r in RESULTS:
            if not r["ok"]:
                print(f"  [{r['path_id']}] {r['name']}")
                if r["notes"]:
                    print(f"        {r['notes']}")


if __name__ == "__main__":
    main()
