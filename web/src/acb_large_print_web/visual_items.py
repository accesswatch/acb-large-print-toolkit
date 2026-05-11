"""Extract visual items from uploaded files for AI-assisted alt-text workflows."""

from __future__ import annotations

import base64
import io
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET


_IMAGE_EXTS = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
    ".tiff": "image/tiff",
}

_WORD_NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _detect_image_dimensions(image_bytes: bytes) -> tuple[int | None, int | None]:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(image_bytes)) as img:
            width, height = img.size
            return int(width), int(height)
    except Exception:
        return None, None


def _mime_for_ext(ext: str) -> str:
    return _IMAGE_EXTS.get(ext.lower(), "image/png")


def _make_preview_data_uri(image_bytes: bytes, mime_type: str) -> str:
    encoded = base64.b64encode(image_bytes).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def _safe_text(value: object, limit: int = 500) -> str:
    text = str(value or "").strip()
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "..."


def _context_lines(*items: object) -> list[str]:
    lines: list[str] = []
    for item in items:
        text = _safe_text(item, 500)
        if text:
            lines.append(text)
    return lines


def _image_item(
    *,
    label: str,
    location: str,
    image_bytes: bytes,
    ext: str,
    current_alt_text: str = "",
    context_lines: list[str] | None = None,
    source_type: str,
) -> dict[str, Any]:
    mime_type = _mime_for_ext(ext)
    width, height = _detect_image_dimensions(image_bytes)
    return {
        "label": label,
        "location": location,
        "source_type": source_type,
        "mime_type": mime_type,
        "image_bytes": image_bytes,
        "preview_data_uri": _make_preview_data_uri(image_bytes, mime_type),
        "width": width,
        "height": height,
        "current_alt_text": _safe_text(current_alt_text, 300),
        "context_lines": context_lines or [],
        "text_only": False,
    }


def _text_item(
    *,
    label: str,
    location: str,
    context_lines: list[str],
    current_alt_text: str = "",
    source_type: str,
) -> dict[str, Any]:
    return {
        "label": label,
        "location": location,
        "source_type": source_type,
        "mime_type": "",
        "image_bytes": b"",
        "preview_data_uri": "",
        "width": None,
        "height": None,
        "current_alt_text": _safe_text(current_alt_text, 300),
        "context_lines": context_lines,
        "text_only": True,
    }


def _extract_from_image_file(path: Path) -> list[dict[str, Any]]:
    ext = path.suffix.lower()
    if ext not in _IMAGE_EXTS:
        return []
    return [
        _image_item(
            label=path.name,
            location="Uploaded image file",
            image_bytes=path.read_bytes(),
            ext=ext,
            context_lines=_context_lines(f"Filename: {path.name}"),
            source_type="image-file",
        )
    ]


def _extract_docx_media(path: Path) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    try:
        with zipfile.ZipFile(path, "r") as archive:
            rels_root = ET.fromstring(archive.read("word/_rels/document.xml.rels"))
            rel_map = {
                rel.attrib.get("Id", ""): rel.attrib.get("Target", "")
                for rel in rels_root
                if rel.attrib.get("Id")
            }
            doc_root = ET.fromstring(archive.read("word/document.xml"))
            doc_texts = [
                _safe_text(node.text, 200)
                for node in doc_root.findall(".//w:t", _WORD_NS)
                if (node.text or "").strip()
            ]
            surrounding = "; ".join(doc_texts[:6])

            for idx, drawing in enumerate(doc_root.findall(".//w:drawing", _WORD_NS), start=1):
                doc_pr = drawing.find(".//wp:docPr", _WORD_NS)
                current_alt = (doc_pr.attrib.get("descr", "") if doc_pr is not None else "").strip()
                blip = drawing.find(".//a:blip", _WORD_NS)
                rel_id = ""
                if blip is not None:
                    rel_id = blip.attrib.get("{%s}embed" % _WORD_NS["r"], "")
                target = rel_map.get(rel_id, "")
                if not target:
                    continue
                normalized = target.replace("..", "").lstrip("/")
                media_name = f"word/{normalized}" if not normalized.startswith("word/") else normalized
                ext = Path(media_name).suffix.lower()
                if ext not in _IMAGE_EXTS:
                    continue
                try:
                    image_bytes = archive.read(media_name)
                except KeyError:
                    continue
                items.append(
                    _image_item(
                        label=f"Word image {idx}",
                        location=f"Word document image {idx}",
                        image_bytes=image_bytes,
                        ext=ext,
                        current_alt_text=current_alt,
                        context_lines=_context_lines(
                            f"Source document: {path.name}",
                            f"Current alt text: {current_alt}" if current_alt else "",
                            f"Document text excerpt: {surrounding}" if surrounding else "",
                        ),
                        source_type="docx-image",
                    )
                )
    except Exception:
        return []
    return items


def _ppt_shape_alt(shape: Any) -> str:
    try:
        c_nv_pr = shape._element.xpath('.//*[local-name()="cNvPr"]')
        if c_nv_pr:
            return str(c_nv_pr[0].get("descr") or "").strip()
    except Exception:
        pass
    return ""


def _ppt_chart_summary(shape: Any) -> str:
    try:
        chart = shape.chart
    except Exception:
        return ""

    bits: list[str] = []
    try:
        if chart.chart_title and chart.chart_title.has_text_frame:
            bits.append(f"Chart title: {chart.chart_title.text_frame.text}")
    except Exception:
        pass
    try:
        plot = chart.plots[0]
        bits.append(f"Chart type: {plot.__class__.__name__}")
    except Exception:
        pass
    try:
        series_bits: list[str] = []
        for series in chart.series[:3]:
            name = getattr(series, "name", "") or "Unnamed series"
            values = []
            try:
                values = [str(v) for v in list(series.values)[:4]]
            except Exception:
                values = []
            if values:
                series_bits.append(f"{name}: {', '.join(values)}")
            else:
                series_bits.append(str(name))
        if series_bits:
            bits.append("Series summary: " + "; ".join(series_bits))
    except Exception:
        pass
    return _safe_text(" ".join(bits), 600)


def _extract_pptx_visuals(path: Path) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:
        return []

    items: list[dict[str, Any]] = []
    try:
        prs = Presentation(str(path))
    except Exception:
        return []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes_text = ""

        shape_counter = 0
        for shape in slide.shapes:
            shape_counter += 1
            current_alt = _ppt_shape_alt(shape)
            base_context = _context_lines(
                f"Presentation: {path.name}",
                f"Slide {slide_index}",
                f"Slide title: {slide_title}" if slide_title else "",
                f"Speaker notes: {notes_text}" if notes_text else "",
                f"Current alt text: {current_alt}" if current_alt else "",
            )
            if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                try:
                    blob = shape.image.blob
                    ext = "." + str(shape.image.ext).lower().lstrip(".")
                    if ext not in _IMAGE_EXTS:
                        ext = ".png"
                    items.append(
                        _image_item(
                            label=f"Slide {slide_index} image {shape_counter}",
                            location=f"Slide {slide_index}",
                            image_bytes=blob,
                            ext=ext,
                            current_alt_text=current_alt,
                            context_lines=base_context,
                            source_type="pptx-image",
                        )
                    )
                except Exception:
                    continue
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart_summary = _ppt_chart_summary(shape)
                items.append(
                    _text_item(
                        label=f"Slide {slide_index} chart {shape_counter}",
                        location=f"Slide {slide_index}",
                        current_alt_text=current_alt,
                        context_lines=base_context + _context_lines(chart_summary),
                        source_type="pptx-chart",
                    )
                )
    return items


def _xlsx_chart_summary(chart: Any) -> str:
    bits: list[str] = []
    title = getattr(chart, "title", None)
    if title:
        bits.append(f"Chart title: {title}")
    bits.append(f"Chart type: {chart.__class__.__name__}")
    series = getattr(chart, "ser", None) or getattr(chart, "series", None) or []
    if series:
        names: list[str] = []
        for item in list(series)[:4]:
            name = getattr(item, "title", None) or getattr(item, "tx", None) or "Series"
            names.append(_safe_text(name, 120))
        bits.append("Series: " + ", ".join(names))
    return _safe_text(" ".join(bits), 600)


def _extract_xlsx_visuals(path: Path) -> list[dict[str, Any]]:
    try:
        from openpyxl import load_workbook
    except Exception:
        return []

    items: list[dict[str, Any]] = []
    try:
        wb = load_workbook(str(path), data_only=True)
    except Exception:
        return []

    for ws in wb.worksheets:
        for idx, img in enumerate(getattr(ws, "_images", []), start=1):
            try:
                image_bytes = img._data() if callable(getattr(img, "_data", None)) else b""
            except Exception:
                image_bytes = b""
            if not image_bytes:
                continue
            ext = Path(getattr(img, "path", "") or "image.png").suffix.lower() or ".png"
            if ext not in _IMAGE_EXTS:
                ext = ".png"
            desc = str(getattr(img, "desc", "") or "").strip()
            anchor = _safe_text(getattr(img, "anchor", ""), 120)
            items.append(
                _image_item(
                    label=f"{ws.title} image {idx}",
                    location=f"Sheet {ws.title}",
                    image_bytes=image_bytes,
                    ext=ext,
                    current_alt_text=desc,
                    context_lines=_context_lines(
                        f"Workbook: {path.name}",
                        f"Worksheet: {ws.title}",
                        f"Anchor: {anchor}" if anchor else "",
                        f"Current alt text: {desc}" if desc else "",
                    ),
                    source_type="xlsx-image",
                )
            )

        for idx, chart in enumerate(getattr(ws, "_charts", []), start=1):
            items.append(
                _text_item(
                    label=f"{ws.title} chart {idx}",
                    location=f"Sheet {ws.title}",
                    context_lines=_context_lines(
                        f"Workbook: {path.name}",
                        f"Worksheet: {ws.title}",
                        _xlsx_chart_summary(chart),
                    ),
                    source_type="xlsx-chart",
                )
            )
    try:
        wb.close()
    except Exception:
        pass
    return items


def _extract_pdf_visuals(path: Path) -> list[dict[str, Any]]:
    try:
        import fitz
    except Exception:
        return []

    items: list[dict[str, Any]] = []
    try:
        doc = fitz.open(str(path))
    except Exception:
        return []

    for page_index in range(min(len(doc), 10)):
        page = doc[page_index]
        page_text = _safe_text(page.get_text("text"), 500)
        page_images = page.get_images(full=True)
        seen_xrefs: set[int] = set()
        for img_index, image_meta in enumerate(page_images, start=1):
            xref = int(image_meta[0])
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                extracted = doc.extract_image(xref)
            except Exception:
                continue
            image_bytes = extracted.get("image") or b""
            ext = "." + str(extracted.get("ext") or "png").lower().lstrip(".")
            if not image_bytes or ext not in _IMAGE_EXTS:
                ext = ".png"
            if image_bytes:
                items.append(
                    _image_item(
                        label=f"PDF page {page_index + 1} image {img_index}",
                        location=f"Page {page_index + 1}",
                        image_bytes=image_bytes,
                        ext=ext,
                        context_lines=_context_lines(
                            f"PDF: {path.name}",
                            f"Page {page_index + 1}",
                            f"Nearby page text: {page_text}" if page_text else "",
                        ),
                        source_type="pdf-image",
                    )
                )
        if not page_images:
            try:
                pix = page.get_pixmap(dpi=140, alpha=False)
                rendered = pix.tobytes("png")
                if rendered:
                    items.append(
                        _image_item(
                            label=f"PDF page {page_index + 1} render",
                            location=f"Page {page_index + 1}",
                            image_bytes=rendered,
                            ext=".png",
                            context_lines=_context_lines(
                                f"PDF: {path.name}",
                                f"Rendered page {page_index + 1}",
                                f"Nearby page text: {page_text}" if page_text else "",
                            ),
                            source_type="pdf-page",
                        )
                    )
            except Exception:
                continue
    doc.close()
    return items


def _extract_epub_visuals(path: Path) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    try:
        with zipfile.ZipFile(path, "r") as archive:
            for idx, name in enumerate(sorted(archive.namelist()), start=1):
                ext = Path(name).suffix.lower()
                if ext not in _IMAGE_EXTS or name.endswith("/"):
                    continue
                try:
                    image_bytes = archive.read(name)
                except KeyError:
                    continue
                items.append(
                    _image_item(
                        label=f"EPUB image {idx}",
                        location=name,
                        image_bytes=image_bytes,
                        ext=ext,
                        context_lines=_context_lines(f"EPUB: {path.name}", f"Archive path: {name}"),
                        source_type="epub-image",
                    )
                )
    except Exception:
        return []
    return items


def extract_visual_items(path: Path, max_items: int = 24) -> list[dict[str, Any]]:
    """Return a list of visual items that can be used for alt-text drafting."""
    path = Path(path)
    ext = path.suffix.lower()

    if ext in _IMAGE_EXTS:
        items = _extract_from_image_file(path)
    elif ext == ".docx":
        items = _extract_docx_media(path)
    elif ext == ".pptx":
        items = _extract_pptx_visuals(path)
    elif ext == ".xlsx":
        items = _extract_xlsx_visuals(path)
    elif ext == ".pdf":
        items = _extract_pdf_visuals(path)
    elif ext == ".epub":
        items = _extract_epub_visuals(path)
    else:
        items = []

    trimmed = items[:max_items]
    for idx, item in enumerate(trimmed):
        item["index"] = idx
        item["total_items"] = len(trimmed)
    return trimmed
