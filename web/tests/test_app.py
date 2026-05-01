"""Tests for the ACB Large Print web application."""

from __future__ import annotations

import io
import json
import re
import sqlite3
import sys
import time
import types
import wave
import zipfile
from contextlib import nullcontext
from pathlib import Path

import pytest
from flask import Flask
from acb_large_print_web import ai_features

from acb_large_print.stress_profiles import describe_stress_corpus
from acb_large_print_web.app import create_app


@pytest.fixture()
def app(tmp_path: Path) -> Flask:
    """Create a test Flask application with a temporary instance folder."""
    application = create_app(
        {
            "TESTING": True,
            "WTF_CSRF_ENABLED": False,
            "MAX_CONTENT_LENGTH": 500 * 1024 * 1024,
        }
    )
    application.instance_path = str(tmp_path / "instance")
    Path(application.instance_path).mkdir(parents=True, exist_ok=True)
    return application


@pytest.fixture()
def client(app: Flask):
    """Flask test client."""
    return app.test_client()


def _make_fake_docx() -> io.BytesIO:
    """Create a minimal valid .docx file (ZIP with correct structure)."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        # Minimal [Content_Types].xml
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>",
        )
        # Minimal relationships
        zf.writestr(
            "_rels/.rels",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
            'Target="word/document.xml"/>'
            "</Relationships>",
        )
        # Minimal document.xml
        zf.writestr(
            "word/document.xml",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body>"
            "<w:p><w:r><w:t>Test</w:t></w:r></w:p>"
            "</w:body>"
            "</w:document>",
        )
    buf.seek(0)
    return buf


def _make_fake_xlsx() -> io.BytesIO:
    """Create a minimal valid .xlsx file using openpyxl."""
    from openpyxl import Workbook

    buf = io.BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Test"
    wb.save(buf)
    buf.seek(0)
    return buf


def _make_fake_pptx() -> io.BytesIO:
    """Create a minimal valid .pptx file using python-pptx."""
    from pptx import Presentation

    buf = io.BytesIO()
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    prs.save(buf)
    buf.seek(0)
    return buf


# ============================================================
# Smoke tests -- every page loads
# ============================================================


class TestPageLoads:
    """Every GET endpoint should return 200."""

    def test_home(self, client):
        resp = client.get("/")
        assert resp.status_code == 200
        assert b"GLOW Accessibility Toolkit" in resp.data

    def test_home_footer_shows_status_link(self, client):
        resp = client.get("/")
        assert resp.status_code == 200
        assert b"/status" in resp.data
        assert b"Server Status" in resp.data

    def test_anthem_download_route(self, client):
        resp = client.get("/anthem/download")
        assert resp.status_code == 200
        assert resp.headers.get("Content-Type", "").startswith("audio/mpeg")
        assert "attachment;" in (resp.headers.get("Content-Disposition") or "")

    def test_audit_form(self, client):
        resp = client.get("/audit/")
        assert resp.status_code == 200
        assert b"Audit" in resp.data

    def test_fix_form(self, client):
        resp = client.get("/fix/")
        assert resp.status_code == 200
        assert b"Fix" in resp.data

    def test_template_form(self, client):
        resp = client.get("/template/")
        assert resp.status_code == 200
        assert b"Template" in resp.data

    def test_export_form(self, client):
        resp = client.get("/export/")
        assert resp.status_code == 301
        assert "/convert/" in resp.headers.get("Location", "")

    def test_guidelines(self, client):
        resp = client.get("/guidelines/")
        assert resp.status_code == 200
        assert b"Guidelines" in resp.data

    def test_feedback_form(self, client):
        resp = client.get("/feedback/")
        assert resp.status_code == 200
        assert b"Feedback" in resp.data

    def test_faq_page(self, client):
        resp = client.get("/faq/")
        assert resp.status_code == 200
        assert b"Frequently Asked Questions" in resp.data

    def test_settings_page(self, client):
        resp = client.get("/settings/")
        assert resp.status_code == 200
        assert b"Settings" in resp.data

    def test_guide_mentions_stress_testing(self, client):
        resp = client.get("/guide/")
        assert resp.status_code == 200
        assert b"Stress Testing and Product Learning" in resp.data
        # Guide is now static markdown; the corpus count is written as "1,000" (with comma)
        total = describe_stress_corpus()["total_heading_cases"]
        assert f"{total:,}".encode() in resp.data

    def test_about_mentions_stress_harness(self, client):
        resp = client.get("/about/")
        assert resp.status_code == 200
        assert b"Stress Testing" in resp.data
        assert b"1000" in resp.data

    def test_health(self, client):
        resp = client.get("/health")
        assert resp.status_code == 200
        payload = resp.get_json()
        assert payload is not None
        assert payload["status"] in {"ok", "degraded"}
        assert "services" in payload
        assert {"web", "openrouter", "whisper", "speech", "braille"}.issubset(set(payload["services"].keys()))
        assert "feature_flags" in payload
        assert "readiness" in payload
        assert {"speech", "braille"}.issubset(set(payload["readiness"].keys()))

    def test_status_page(self, client):
        resp = client.get("/status")
        assert resp.status_code == 200
        body = resp.get_data(as_text=True)
        assert "Server Status" in body
        assert "Raw Health JSON" in body
        assert "Service Reachability" in body

    def test_status_page_handles_non_string_louis_version(self, client, monkeypatch):
        import acb_large_print.braille_converter as braille_converter

        class _FakeVersion:
            def __repr__(self):
                return "mock-version"

        monkeypatch.setattr(braille_converter, "louis_version", lambda: _FakeVersion())
        monkeypatch.setattr(braille_converter, "braille_available", lambda: True)

        resp = client.get("/status")
        assert resp.status_code == 200
        body = resp.get_data(as_text=True)
        assert "Server Status" in body


@pytest.mark.skipif(
    not ai_features.ai_whisperer_enabled(),
    reason="AI whisperer disabled via feature flags; skipping Whisperer tests",
)
class TestWhispererProgress:
    def test_whisperer_estimate_uses_metadata_when_available(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "_estimate_audio_duration_seconds", lambda _p: 120.0)

        resp = client.post(
            "/whisperer/estimate",
            # Use a realistic compressed size so metadata remains plausible.
            data={"audio": (io.BytesIO(b"x" * (3 * 1024 * 1024)), "meeting.mp3")},
            content_type="multipart/form-data",
        )

        assert resp.status_code == 200
        payload = resp.get_json()
        assert payload is not None
        assert payload["source"] == "metadata"
        assert payload["audio_seconds"] == 120.0
        assert payload["expected_seconds"] >= 120.0

    def test_whisperer_estimate_falls_back_to_size(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "_estimate_audio_duration_seconds", lambda _p: None)

        # 16000 bytes ~= 1 second using the fallback heuristic.
        resp = client.post(
            "/whisperer/estimate",
            data={"audio": (io.BytesIO(b"x" * 16000), "meeting.mp3")},
            content_type="multipart/form-data",
        )

        assert resp.status_code == 200
        payload = resp.get_json()
        assert payload is not None
        assert payload["source"] == "size-fallback"
        assert payload["audio_seconds"] >= 1.0
        assert payload["expected_seconds"] >= 15.0

    def test_whisperer_estimate_page_preserves_token_for_start(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "is_whisper_configured", lambda: True)
        monkeypatch.setattr(whisperer_route, "audio_gate", lambda: nullcontext())
        monkeypatch.setattr(whisperer_route, "_estimate_audio_duration_seconds", lambda _p: 120.0)

        def _fake_transcribe(_audio_path, language=None, session_hash=""):
            assert session_hash
            return "# Token Test Transcript\n"

        monkeypatch.setattr(whisperer_route, "gateway_transcribe", _fake_transcribe)

        estimate_page = client.post(
            "/whisperer/estimate-page",
            data={"audio": (io.BytesIO(b"x" * (2 * 1024 * 1024)), "meeting.mp3")},
            content_type="multipart/form-data",
        )
        assert estimate_page.status_code == 200
        html = estimate_page.get_data(as_text=True)

        import re

        token_match = re.search(r'name="existing_token" value="([^"]+)"', html)
        assert token_match is not None
        token = token_match.group(1)

        start = client.post(
            "/whisperer/start",
            data={
                "existing_token": token,
                "uploaded_filename": "meeting.mp3",
                "output_format": "markdown",
                "confirm_estimate": "yes",
                "estimate_source": "metadata",
            },
            content_type="multipart/form-data",
        )
        assert start.status_code == 202
        payload = start.get_json()
        assert payload is not None
        job_id = payload["job_id"]

        final_payload = None
        for _ in range(20):
            prog = client.get(f"/whisperer/progress/{job_id}")
            assert prog.status_code == 200
            final_payload = prog.get_json()
            if final_payload and final_payload.get("status") == "complete":
                break
            time.sleep(0.02)

        assert final_payload is not None
        assert final_payload["status"] == "complete"
        assert final_payload["progress"] == 100

    def test_whisperer_background_progress_and_download(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "is_whisper_configured", lambda: True)
        monkeypatch.setattr(whisperer_route, "audio_gate", lambda: nullcontext())

        def _fake_transcribe(_audio_path, language=None, session_hash=""):
            return "# Test Transcript\n\nHello world.\n"

        monkeypatch.setattr(whisperer_route, "gateway_transcribe", _fake_transcribe)

        start = client.post(
            "/whisperer/start",
            data={
                "audio": (io.BytesIO(b"fake audio bytes"), "meeting.mp3"),
                "output_format": "markdown",
                "language": "en",
                "confirm_estimate": "yes",
            },
            content_type="multipart/form-data",
        )
        assert start.status_code == 202
        payload = start.get_json()
        assert payload is not None
        job_id = payload["job_id"]

        final_payload = None
        for _ in range(20):
            prog = client.get(f"/whisperer/progress/{job_id}")
            assert prog.status_code == 200
            final_payload = prog.get_json()
            if final_payload and final_payload.get("status") == "complete":
                break
            time.sleep(0.02)

        assert final_payload is not None
        assert final_payload["status"] == "complete"
        assert final_payload["progress"] == 100
        assert "download_url" in final_payload

        download = client.get(final_payload["download_url"])
        assert download.status_code == 200
        assert b"Test Transcript" in download.data

    def test_whisperer_requires_estimate_ack(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "is_whisper_configured", lambda: True)

        resp = client.post(
            "/whisperer/start",
            data={
                "audio": (io.BytesIO(b"fake audio bytes"), "meeting.mp3"),
                "output_format": "markdown",
                # Intentionally omit confirm_estimate
            },
            content_type="multipart/form-data",
        )

        assert resp.status_code == 400
        payload = resp.get_json()
        assert payload is not None
        assert "check the confirmation box" in payload["error"]

    def test_whisperer_rejects_audio_above_duration_limit(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "is_whisper_configured", lambda: True)
        monkeypatch.setattr(whisperer_route, "_estimate_audio_duration_seconds", lambda _p: 3 * 60 * 60)
        monkeypatch.setattr(whisperer_route, "_sanitize_duration_estimate", lambda _p, d: d)
        monkeypatch.setattr(whisperer_route, "_MAX_AUDIO_MINUTES", 120)

        resp = client.post(
            "/whisperer/start",
            data={
                "audio": (io.BytesIO(b"fake audio bytes"), "meeting.mp3"),
                "output_format": "markdown",
                "confirm_estimate": "yes",
            },
            content_type="multipart/form-data",
        )

        assert resp.status_code == 400
        payload = resp.get_json()
        assert payload is not None
        assert "maximum supported length" in payload["error"]

    def test_background_opt_in_requires_email_service(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "is_whisper_configured", lambda: True)
        monkeypatch.setattr(whisperer_route, "email_configured", lambda: False)
        monkeypatch.setattr(whisperer_route, "_estimate_audio_duration_seconds", lambda _p: 45 * 60)

        resp = client.post(
            "/whisperer/start",
            data={
                "audio": (io.BytesIO(b"fake audio bytes"), "meeting.mp3"),
                "output_format": "markdown",
                "confirm_estimate": "yes",
                "background_opt_in": "yes",
                "notify_email": "user@example.com",
                "retrieval_password": "secret123!",
                "retrieval_password_confirm": "secret123!",
            },
            content_type="multipart/form-data",
        )

        assert resp.status_code == 400
        payload = resp.get_json()
        assert payload is not None
        assert "requires email service" in payload["error"]

    def test_background_secure_retrieve_requires_password(self, client, monkeypatch):
        import acb_large_print_web.routes.whisperer as whisperer_route

        monkeypatch.setattr(whisperer_route, "is_whisper_configured", lambda: True)
        monkeypatch.setattr(whisperer_route, "email_configured", lambda: True)
        monkeypatch.setattr(whisperer_route, "send_whisperer_status_email", lambda *a, **k: (True, "ok"))
        monkeypatch.setattr(whisperer_route, "audio_gate", lambda: nullcontext())
        monkeypatch.setattr(whisperer_route, "_estimate_audio_duration_seconds", lambda _p: 45 * 60)

        def _fake_transcribe(_audio_path, language=None, session_hash=""):
            return "# Secure Transcript\n\nHello world.\n"

        monkeypatch.setattr(whisperer_route, "gateway_transcribe", _fake_transcribe)

        start = client.post(
            "/whisperer/start",
            data={
                "audio": (io.BytesIO(b"fake audio bytes"), "meeting.mp3"),
                "output_format": "markdown",
                "confirm_estimate": "yes",
                "background_opt_in": "yes",
                "notify_email": "user@example.com",
                "retrieval_password": "secret123!",
                "retrieval_password_confirm": "secret123!",
            },
            content_type="multipart/form-data",
        )
        assert start.status_code == 202
        job_id = start.get_json()["job_id"]

        # Wait for completion
        final_payload = None
        for _ in range(20):
            prog = client.get(f"/whisperer/progress/{job_id}")
            assert prog.status_code == 200
            final_payload = prog.get_json()
            if final_payload and final_payload.get("status") == "complete":
                break
            time.sleep(0.02)

        assert final_payload is not None
        assert final_payload["status"] == "complete"
        assert final_payload["background_opt_in"] is True

        # Find retrieval token from in-memory job registry
        job = whisperer_route._get_job(job_id)
        assert job is not None
        token = job.retrieval_token
        assert token

        bad = client.post(f"/whisperer/retrieve/{token}", data={"retrieval_password": "wrong"})
        assert bad.status_code == 403

        good = client.post(f"/whisperer/retrieve/{token}", data={"retrieval_password": "secret123!"})
        assert good.status_code == 200
        assert b"Secure Transcript" in good.data


class TestAdminAuth:
    def test_admin_login_page_loads(self, client, monkeypatch):
        import acb_large_print_web.routes.admin as admin_route

        monkeypatch.setattr(admin_route, "email_configured", lambda: True)
        resp = client.get("/admin/login")
        assert resp.status_code == 200
        assert b"Admin Sign-In" in resp.data
        html = resp.data.decode("utf-8")
        # Check for accessible email input and send button (requires email_configured=True)
        assert 'id="email"' in html
        assert "Admin email" in html
        assert 'id="admin-login-btn"' in html

    def test_admin_request_access_requires_email_config(self, client):
        resp = client.get("/admin/request-access")
        assert resp.status_code == 503
        assert b"unavailable until email delivery is configured" in resp.data
        html = resp.data.decode("utf-8")
        # When email not configured, the request-access form should not be shown
        assert "Email" not in html
        assert 'id="access-request-btn"' not in html

    def test_admin_magic_link_sign_in_for_bootstrap_admin(self, app, client, monkeypatch):
        import acb_large_print_web.routes.admin as admin_route

        monkeypatch.setenv("ADMIN_BOOTSTRAP_EMAILS", "admin@example.com")
        monkeypatch.setattr(admin_route, "email_configured", lambda: True)
        monkeypatch.setattr(admin_route, "_send_admin_email", lambda *a, **k: None)
        monkeypatch.setattr(admin_route.secrets, "token_urlsafe", lambda _n: "admintoken")

        send_resp = client.post(
            "/admin/login/email",
            data={"email": "admin@example.com"},
        )
        assert send_resp.status_code == 200
        assert b"sign-in link" in send_resp.data

        consume_resp = client.get("/admin/magic-link/consume?token=admintoken")
        assert consume_resp.status_code == 302
        assert "/admin/queue" in consume_resp.location

        queue_resp = client.get("/admin/queue")
        assert queue_resp.status_code == 200
        assert b"Admin Queue Dashboard" in queue_resp.data

    def test_admin_speech_page_loads_for_signed_in_admin(self, client, monkeypatch):
        import acb_large_print_web.routes.admin as admin_route

        monkeypatch.setenv("ADMIN_BOOTSTRAP_EMAILS", "admin@example.com")
        monkeypatch.setattr(admin_route, "email_configured", lambda: True)
        monkeypatch.setattr(admin_route, "_send_admin_email", lambda *a, **k: None)
        monkeypatch.setattr(admin_route.secrets, "token_urlsafe", lambda _n: "admintoken2")
        monkeypatch.setattr(
            admin_route,
            "get_engine_status",
            lambda: {
                "kokoro": {
                    "installed": True,
                    "models_present": True,
                    "ready": True,
                    "voices_available": ["af_bella"],
                    "model_dir": "/tmp/speech_models",
                    "setup_commands": [],
                },
                "piper": {
                    "installed": True,
                    "models_present": True,
                    "ready": True,
                    "voices_available": ["en_US-lessac-medium"],
                    "model_dir": "/tmp/speech_models/piper",
                    "setup_commands": [],
                },
            },
        )
        monkeypatch.setattr(
            admin_route,
            "get_piper_voice_inventory",
            lambda: [
                {
                    "id": "en_US-lessac-medium",
                    "label": "Lessac (US)",
                    "accent": "American",
                    "gender": "Male",
                    "sample_rate": 22050,
                    "installed": True,
                    "model_path": "/tmp/en_US-lessac-medium.onnx",
                    "config_path": "/tmp/en_US-lessac-medium.onnx.json",
                }
            ],
        )

        client.post("/admin/login/email", data={"email": "admin@example.com"})
        client.get("/admin/magic-link/consume?token=admintoken2")

        resp = client.get("/admin/speech")
        assert resp.status_code == 200
        assert b"Admin Speech Studio" in resp.data
        assert b"en_US-lessac-medium" in resp.data

    def test_admin_speech_install_action_redirects_with_notice(self, client, monkeypatch):
        import acb_large_print_web.routes.admin as admin_route

        monkeypatch.setenv("ADMIN_BOOTSTRAP_EMAILS", "admin@example.com")
        monkeypatch.setattr(admin_route, "email_configured", lambda: True)
        monkeypatch.setattr(admin_route, "_send_admin_email", lambda *a, **k: None)
        monkeypatch.setattr(admin_route.secrets, "token_urlsafe", lambda _n: "admintoken3")
        monkeypatch.setattr(admin_route, "install_piper_voice", lambda _voice: (True, "Installed test voice."))

        client.post("/admin/login/email", data={"email": "admin@example.com"})
        client.get("/admin/magic-link/consume?token=admintoken3")

        resp = client.post("/admin/speech/install/en_US-lessac-medium")
        assert resp.status_code == 302
        assert "/admin/speech" in resp.location
        assert "notice=Installed+test+voice." in resp.location


# ============================================================
# Error handling
# ============================================================


class TestErrors:
    def test_404(self, client):
        resp = client.get("/nonexistent-page/")
        assert resp.status_code == 404
        assert b"Page Not Found" in resp.data

    def test_audit_no_file(self, client):
        resp = client.post("/audit/", data={})
        assert resp.status_code == 400
        assert b"No file selected" in resp.data

    def test_fix_no_file(self, client):
        resp = client.post("/fix/", data={})
        assert resp.status_code == 400
        assert b"No file selected" in resp.data

    def test_audit_wrong_extension(self, client):
        data = {"document": (io.BytesIO(b"not a docx"), "test.txt")}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"not supported" in resp.data

    def test_fix_wrong_extension(self, client):
        data = {"document": (io.BytesIO(b"not a docx"), "test.txt")}
        resp = client.post("/fix/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"not supported" in resp.data

    def test_audit_corrupt_file(self, client):
        """A .docx that isn't a valid ZIP should be rejected."""
        data = {"document": (io.BytesIO(b"not a zip archive"), "test.docx")}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"does not appear to be a valid" in resp.data


# ============================================================
# Feedback
# ============================================================


class TestFeedback:
    def test_submit_feedback(self, app, client):
        resp = client.post(
            "/feedback/",
            data={
                "rating": "good",
                "task": "audit",
                "message": "Works great!",
            },
        )
        assert resp.status_code == 200
        assert b"Thank You" in resp.data

        # Verify it was saved in SQLite
        db_path = Path(app.instance_path) / "feedback.db"
        assert db_path.exists()
        conn = sqlite3.connect(str(db_path))
        rows = conn.execute("SELECT rating, message FROM feedback").fetchall()
        conn.close()
        assert len(rows) == 1
        assert rows[0][0] == "good"
        assert rows[0][1] == "Works great!"

    def test_submit_feedback_missing_rating(self, client):
        resp = client.post(
            "/feedback/",
            data={
                "message": "Some feedback",
            },
        )
        assert resp.status_code == 400
        assert b"select a rating" in resp.data

    def test_submit_feedback_missing_message(self, client):
        resp = client.post(
            "/feedback/",
            data={
                "rating": "excellent",
            },
        )
        assert resp.status_code == 400
        assert b"enter your feedback" in resp.data

    def test_review_no_password_configured(self, client):
        """Review returns 404 when FEEDBACK_PASSWORD is not set."""
        resp = client.get("/feedback/review")
        assert resp.status_code == 404

    def test_review_wrong_password(self, app, client):
        app.config["FEEDBACK_PASSWORD"] = "secret123"
        import os

        os.environ["FEEDBACK_PASSWORD"] = "secret123"
        resp = client.get("/feedback/review?key=wrong")
        assert resp.status_code == 403
        os.environ.pop("FEEDBACK_PASSWORD", None)

    def test_review_correct_password(self, app, client):
        import os

        os.environ["FEEDBACK_PASSWORD"] = "secret123"
        # Submit a feedback entry first
        client.post(
            "/feedback/",
            data={
                "rating": "excellent",
                "task": "fix",
                "message": "Love it",
            },
        )
        resp = client.get("/feedback/review?key=secret123")
        assert resp.status_code == 200
        assert b"Love it" in resp.data
        assert b"Excellent" in resp.data
        os.environ.pop("FEEDBACK_PASSWORD", None)


# ============================================================
# Accessibility checks on rendered HTML
# ============================================================


class TestAccessibility:
    """Basic structural accessibility checks on rendered pages."""

    def test_all_pages_have_lang(self, client):
        for path in [
            "/",
            "/audit/",
            "/fix/",
            "/faq/",
            "/guidelines/",
            "/feedback/",
            "/about/",
            "/convert/",
        ]:
            resp = client.get(path)
            assert b'lang="en"' in resp.data, f"Missing lang on {path}"

    def test_all_pages_have_main_landmark(self, client):
        for path in [
            "/",
            "/audit/",
            "/fix/",
            "/faq/",
            "/guidelines/",
            "/feedback/",
            "/about/",
            "/convert/",
        ]:
            resp = client.get(path)
            assert b'id="main"' in resp.data, f"Missing main landmark on {path}"

    def test_nav_has_aria_label(self, client):
        resp = client.get("/")
        assert b"aria-label" in resp.data

    def test_footer_has_role(self, client):
        resp = client.get("/")
        assert b"<footer" in resp.data
        assert b'role="contentinfo"' not in resp.data

    def test_bits_renamed_to_solutions(self, client):
        resp = client.get("/")
        assert b"Solutions" in resp.data
        assert b"Specialists" not in resp.data


# ============================================================
# Upload + audit integration (with real docx)
# ============================================================


class TestAuditIntegration:
    def test_audit_full_mode(self, client):
        docx_data = _make_fake_docx()
        data = {
            "document": (docx_data, "test.docx"),
            "mode": "full",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Audit Report" in resp.data or b"Score" in resp.data

    def test_audit_quick_mode(self, client):
        docx_data = _make_fake_docx()
        data = {
            "document": (docx_data, "test.docx"),
            "mode": "quick",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_audit_form_has_standards_profile_controls(self, client):
        resp = client.get("/audit/")
        assert resp.status_code == 200
        assert b'name="standards_profile"' in resp.data
        assert b'value="acb_2025"' in resp.data
        assert b'value="aph_submission"' in resp.data
        assert b'value="combined_strict"' in resp.data

    def test_audit_profile_label_shows_in_report(self, client):
        docx_data = _make_fake_docx()
        data = {
            "document": (docx_data, "test.docx"),
            "mode": "full",
            "standards_profile": "aph_submission",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Standards profile" in resp.data
        assert b"APH Submission" in resp.data


# ============================================================
# Template generation
# ============================================================


class TestTemplateGeneration:
    def test_generate_template(self, client):
        resp = client.post(
            "/template/",
            data={
                "title": "Test Template",
            },
        )
        # Should return 200 with a result page or a download
        assert resp.status_code == 200

    def test_template_form_has_standards_profile_controls(self, client):
        resp = client.get("/template/")
        assert resp.status_code == 200
        assert b'name="standards_profile"' in resp.data
        assert b'value="acb_2025"' in resp.data
        assert b'value="aph_submission"' in resp.data
        assert b'value="combined_strict"' in resp.data
        assert b'name="allowed_heading_levels"' in resp.data

    def test_generate_template_aph_profile(self, client):
        resp = client.post(
            "/template/",
            data={
                "title": "APH Template",
                "standards_profile": "aph_submission",
            },
        )
        assert resp.status_code == 200

    def test_generate_template_with_restricted_heading_levels(self, client):
        resp = client.post(
            "/template/",
            data={
                "title": "Restricted Levels",
                "include_sample": "on",
                "allowed_heading_levels": ["1", "2"],
            },
        )
        assert resp.status_code == 200


# ============================================================
# Multi-format audit (xlsx, pptx)
# ============================================================


class TestMultiFormatAudit:
    """Upload Excel and PowerPoint files and verify audit works."""

    def test_audit_xlsx(self, client):
        xlsx_data = _make_fake_xlsx()
        data = {
            "document": (xlsx_data, "test.xlsx"),
            "mode": "full",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Audit Report" in resp.data or b"Score" in resp.data

    def test_audit_pptx(self, client):
        pptx_data = _make_fake_pptx()
        data = {
            "document": (pptx_data, "test.pptx"),
            "mode": "full",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Audit Report" in resp.data or b"Score" in resp.data

    def test_audit_xlsx_quick_mode(self, client):
        xlsx_data = _make_fake_xlsx()
        data = {
            "document": (xlsx_data, "test.xlsx"),
            "mode": "quick",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_audit_pptx_quick_mode(self, client):
        pptx_data = _make_fake_pptx()
        data = {
            "document": (pptx_data, "test.pptx"),
            "mode": "quick",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_corrupt_xlsx_rejected(self, client):
        data = {"document": (io.BytesIO(b"not a zip archive"), "test.xlsx")}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"does not appear to be a valid" in resp.data

    def test_corrupt_pptx_rejected(self, client):
        data = {"document": (io.BytesIO(b"not a zip archive"), "test.pptx")}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"does not appear to be a valid" in resp.data


# ============================================================
# Accessibility structure (extended)
# ============================================================


class TestAccessibilityExtended:
    """Verify WCAG 2.2 AA structural requirements added in multi-format update."""

    def test_skip_link_present(self, client):
        resp = client.get("/")
        assert b"skip-link" in resp.data
        assert b"Skip to main content" in resp.data

    def test_home_shows_all_three_formats(self, client):
        resp = client.get("/")
        assert b"Word (.docx)" in resp.data
        assert b"Excel (.xlsx)" in resp.data
        assert b"PowerPoint (.pptx)" in resp.data

    def test_home_has_settings_tile_and_navigation_descriptions(self, client):
        resp = client.get("/")
        assert b"Settings" in resp.data
        assert b"Main Navigation Tab Order and Purpose" in resp.data
        assert b"Guidelines" in resp.data

    def test_audit_form_accepts_all_formats(self, client):
        resp = client.get("/audit/")
        assert b".docx,.xlsx,.pptx" in resp.data

    def test_flash_error_has_text_prefix(self, client):
        """Error flash messages include 'Error:' text so meaning is not color-alone."""
        resp = client.post("/audit/", data={})
        assert b"<strong>Error:</strong>" in resp.data

    def test_format_tags_in_rule_list(self, client):
        """Rule list checkboxes include format tags with text labels."""
        resp = client.get("/audit/")
        assert b'class="format-tag format-docx"' in resp.data

    def test_severity_badges_have_text(self, client):
        """Severity badges use text labels, not just color."""
        resp = client.get("/audit/")
        for level in [b"Critical", b"High", b"Medium", b"Low"]:
            assert level in resp.data


# ============================================================
# About page
# ============================================================


class TestAboutPage:
    """Tests for the /about/ route."""

    def test_about_page_loads(self, client):
        resp = client.get("/about/")
        assert resp.status_code == 200
        assert b"About" in resp.data

    def test_about_has_lang(self, client):
        resp = client.get("/about/")
        assert b'lang="en"' in resp.data

    def test_about_has_main_landmark(self, client):
        resp = client.get("/about/")
        assert b'id="main"' in resp.data

    def test_about_mentions_markitdown(self, client):
        resp = client.get("/about/")
        assert b"MarkItDown" in resp.data

    def test_about_shows_speech_and_anthem_usage_sections(self, app, client):
        from acb_large_print_web.tool_usage import record as record_usage, record_details

        with app.app_context():
            record_usage("anthem_download")
            record_details(
                "speech",
                {
                    "mode": "typed_preview",
                    "voice": "kokoro:af_bella",
                    "speed": "1.0",
                    "pitch": "0",
                },
            )

        resp = client.get("/about/")
        assert resp.status_code == 200
        assert b"Let it GLOW theme downloads" in resp.data
        assert b"Speech Studio Usage Patterns" in resp.data
        assert b"kokoro:af_bella" in resp.data


class TestSpeechMetrics:
    def test_speech_metrics_records_and_summarizes(self, app):
        from acb_large_print_web import speech_metrics

        with app.app_context():
            speech_metrics.record_document_conversion(
                engine="kokoro",
                voice="kokoro:af_bella",
                speed=1.0,
                pitch=0,
                word_count=120,
                char_count=680,
                source_size_bytes=4096,
                processing_seconds=18.5,
                audio_seconds=42.0,
            )
            summary = speech_metrics.get_summary()

        assert summary["samples"] >= 1
        assert summary["avg_processing_seconds"] > 0
        assert summary["avg_words"] > 0

    def test_speech_metrics_estimator_uses_baseline_when_samples_low(self, app):
        from acb_large_print_web import speech_metrics

        with app.app_context():
            est, source, samples = speech_metrics.estimate_processing_seconds(
                engine="kokoro",
                speed=1.0,
                word_count=200,
                char_count=1200,
                source_size_bytes=8192,
                baseline_seconds=60.0,
            )

        assert est >= 1.0
        assert source in {"baseline", "historical_blended"}
        assert samples >= 0


# ============================================================
# Convert route
# ============================================================


class TestConvertPage:
    """Tests for the /convert/ route."""

    def test_convert_form_loads(self, client):
        resp = client.get("/convert/")
        assert resp.status_code == 200
        assert b"Convert" in resp.data

    def test_convert_has_lang(self, client):
        resp = client.get("/convert/")
        assert b'lang="en"' in resp.data

    def test_convert_has_main_landmark(self, client):
        resp = client.get("/convert/")
        assert b'id="main"' in resp.data

    def test_convert_no_file(self, client):
        resp = client.post("/convert/", data={})
        assert resp.status_code == 400
        assert b"No file selected" in resp.data

    def test_convert_wrong_extension(self, client):
        data = {"document": (io.BytesIO(b"text"), "test.txt")}
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"not supported" in resp.data

    def test_convert_docx_returns_markdown(self, client):
        docx_data = _make_fake_docx()
        data = {"document": (docx_data, "test.docx")}
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        # Route now returns a result page (not a direct file download)
        assert resp.content_type.startswith("text/html")
        assert b"test.md" in resp.data
        assert b"Conversion Complete" in resp.data

    def test_convert_xlsx_returns_markdown(self, client):
        xlsx_data = _make_fake_xlsx()
        data = {"document": (xlsx_data, "test.xlsx")}
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert resp.content_type.startswith("text/html")
        assert b"test.md" in resp.data
        assert b"Conversion Complete" in resp.data

    def test_convert_form_has_direction_radios(self, client):
        resp = client.get("/convert/")
        assert b'name="direction"' in resp.data
        assert b'value="to-markdown"' in resp.data
        assert b'value="to-html"' in resp.data

    def test_convert_form_has_to_speech_direction(self, client):
        resp = client.get("/convert/")
        assert resp.status_code == 200
        assert b'value="to-speech"' in resp.data

    def test_convert_to_speech_redirects_with_token_handoff(self, client):
        data = {
            "document": (_make_fake_docx(), "speech-handoff.docx"),
            "direction": "to-speech",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 302
        assert "/speech/?token=" in (resp.headers.get("Location") or "")

    def test_convert_to_html_md_file(self, client):
        """Upload a .md with direction=to-html; expect HTML back if Pandoc installed."""
        from acb_large_print.pandoc_converter import pandoc_available

        md_content = b"# Hello World\n\nSome text.\n"
        data = {
            "document": (io.BytesIO(md_content), "test.md"),
            "direction": "to-html",
            "acb_format": "on",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        if pandoc_available():
            assert resp.status_code == 200
            assert resp.content_type.startswith("text/html")
            # Route now returns a result page (not a direct file download)
            assert b"test.html" in resp.data
            assert b"Conversion Complete" in resp.data
        else:
            # Pandoc not installed -- should get an error page
            assert resp.status_code in (400, 500)
            assert b"Pandoc" in resp.data

    def test_convert_to_html_unsupported_ext(self, client):
        """Upload a .jpg (image) with direction=to-html -- not a Pandoc-supported format.

        Note: .xlsx was previously used here but is now supported via the two-stage
        MarkItDown→Pandoc chain added in 2.7.0. Images remain unsupported for to-html.
        """
        from acb_large_print.pandoc_converter import pandoc_available

        if not pandoc_available():
            pytest.skip("Pandoc not installed")
        data = {
            "document": (io.BytesIO(b"fake image data"), "test.jpg"),
            "direction": "to-html",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"cannot be converted to HTML" in resp.data

    def test_convert_rst_accepted(self, client):
        """The .rst extension should be accepted by the convert upload validator."""
        data = {
            "document": (io.BytesIO(b"Hello\n=====\n\nSome text.\n"), "test.rst"),
            "direction": "to-markdown",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        # May fail at conversion (MarkItDown may not handle .rst) but should not
        # fail at upload validation (400 "not supported")
        assert resp.status_code != 400 or b"not supported" not in resp.data


# ============================================================
# Markdown and PDF audit
# ============================================================


def _make_fake_md() -> io.BytesIO:
    """Create a simple Markdown file with known issues."""
    content = (
        "# My Document\n\n"
        "Click [here](https://example.com) for details.\n\n"
        "### Skipped heading\n\n"
        "Some *italic* text.\n"
    )
    return io.BytesIO(content.encode("utf-8"))


def _make_fake_pdf() -> io.BytesIO:
    """Create a minimal PDF file."""
    # Minimal valid PDF (1 blank page)
    pdf_bytes = (
        b"%PDF-1.0\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R>>endobj\n"
        b"xref\n0 4\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"trailer<</Size 4/Root 1 0 R>>\n"
        b"startxref\n190\n%%EOF\n"
    )
    return io.BytesIO(pdf_bytes)


class TestMarkdownAudit:
    """Upload Markdown files and verify audit works."""

    def test_audit_markdown_full(self, client):
        md_data = _make_fake_md()
        data = {
            "document": (md_data, "test.md"),
            "mode": "full",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Audit Report" in resp.data or b"Score" in resp.data

    def test_audit_markdown_quick(self, client):
        md_data = _make_fake_md()
        data = {
            "document": (md_data, "test.md"),
            "mode": "quick",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_fix_markdown_returns_warning(self, client):
        """Markdown fix should work but return a 'not yet' warning."""
        md_data = _make_fake_md()
        data = {"document": (md_data, "test.md")}
        resp = client.post("/fix/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200


class TestFixIntegration:
    def test_fix_form_has_standards_profile_controls(self, client):
        resp = client.get("/fix/")
        assert resp.status_code == 200
        assert b'name="standards_profile"' in resp.data
        assert b'value="acb_2025"' in resp.data
        assert b'value="aph_submission"' in resp.data
        assert b'value="combined_strict"' in resp.data

    def test_fix_profile_label_shows_in_results(self, client):
        md_data = _make_fake_md()
        data = {
            "document": (md_data, "test.md"),
            "mode": "full",
            "standards_profile": "combined_strict",
        }
        resp = client.post("/fix/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Standards profile" in resp.data
        assert b"Combined Strict" in resp.data


class TestSettingsIntegration:
    def test_settings_page_has_profile_and_advanced_defaults(self, client):
        resp = client.get("/settings/")
        assert resp.status_code == 200
        assert b'name="settings_audit_profile"' in resp.data
        assert b'name="settings_fix_profile"' in resp.data
        assert b'name="settings_template_profile"' in resp.data
        assert b'name="settings_fix_suppress_link_text"' in resp.data
        assert b'name="settings_fix_use_list_levels"' in resp.data
        assert b'name="settings_fix_list_indent_level_1"' in resp.data
        assert b'name="settings_fix_allowed_heading_levels"' in resp.data
        assert b'name="settings_template_allowed_heading_levels"' in resp.data

    def test_settings_page_shows_speech_studio_by_default(self, client):
        resp = client.get("/settings/")
        assert resp.status_code == 200
        assert b"Speech Studio" in resp.data
        assert b"Open Speech Studio" in resp.data
        assert b"Download sample audio" in resp.data

    def test_speech_setup_references_current_kokoro_model_assets(self):
        from acb_large_print_web import speech

        assert speech._KOKORO_MODEL_FILE == "kokoro-v1.0.onnx"
        assert speech._KOKORO_VOICES_FILE == "voices-v1.0.bin"

    def test_piper_tts_is_installed_with_web_app(self):
        pyproject = Path("web/pyproject.toml").read_text(encoding="utf-8")
        requirements = Path("web/requirements.txt").read_text(encoding="utf-8")

        assert "piper-tts>=1.4.2" in pyproject
        assert "piper-tts>=1.4.2" in requirements

    def test_piper_voice_present_supports_huggingface_nested_layout(self, tmp_path):
        from acb_large_print_web import speech

        old_model_dir = speech._model_dir
        model_root = tmp_path / "speech_models"
        nested = model_root / "piper" / "en" / "en_US" / "lessac" / "medium"
        nested.mkdir(parents=True)
        (nested / "en_US-lessac-medium.onnx").write_bytes(b"model")
        (nested / "en_US-lessac-medium.onnx.json").write_text("{}", encoding="utf-8")

        try:
            speech.configure(model_root)

            assert speech._piper_voice_present("en_US-lessac-medium")
            assert speech._piper_model_path("en_US-lessac-medium") == nested / "en_US-lessac-medium.onnx"
            assert speech._piper_config_path("en_US-lessac-medium") == nested / "en_US-lessac-medium.onnx.json"
        finally:
            speech._model_dir = old_model_dir

    def test_piper_installed_accepts_current_package_import(self, monkeypatch):
        from acb_large_print_web import speech

        fake_piper = types.ModuleType("piper")
        fake_piper.PiperVoice = object
        monkeypatch.setattr("shutil.which", lambda name: None)
        monkeypatch.setitem(sys.modules, "piper", fake_piper)

        assert speech._piper_installed()

    def test_synthesize_piper_uses_current_package_api(self, tmp_path, monkeypatch):
        from acb_large_print_web import speech

        old_model_dir = speech._model_dir
        model_root = tmp_path / "speech_models"
        piper_dir = model_root / "piper"
        piper_dir.mkdir(parents=True)
        (piper_dir / "en_US-lessac-medium.onnx").write_bytes(b"model")
        (piper_dir / "en_US-lessac-medium.onnx.json").write_text("{}", encoding="utf-8")

        calls = {}

        class FakeSynthesisConfig:
            def __init__(self, *, length_scale):
                self.length_scale = length_scale

        class FakeVoice:
            def synthesize_wav(self, text, wav_file, *, syn_config):
                calls["text"] = text
                calls["length_scale"] = syn_config.length_scale
                wav_file.setnchannels(1)
                wav_file.setsampwidth(2)
                wav_file.setframerate(22050)
                wav_file.writeframes(b"\0\0")

        class FakePiperVoice:
            @staticmethod
            def load(model_path, *, config_path=None):
                calls["model_path"] = model_path
                calls["config_path"] = config_path
                return FakeVoice()

        fake_piper = types.ModuleType("piper")
        fake_piper.PiperVoice = FakePiperVoice
        fake_piper.SynthesisConfig = FakeSynthesisConfig
        monkeypatch.setattr("shutil.which", lambda name: None)
        monkeypatch.setitem(sys.modules, "piper", fake_piper)

        try:
            speech.configure(model_root)
            wav_bytes = speech._synthesize_piper("en_US-lessac-medium", "Hello", speed=2.0)
        finally:
            speech._model_dir = old_model_dir

        assert wav_bytes.startswith(b"RIFF")
        assert calls == {
            "model_path": str(piper_dir / "en_US-lessac-medium.onnx"),
            "config_path": str(piper_dir / "en_US-lessac-medium.onnx.json"),
            "text": "Hello",
            "length_scale": 0.5,
        }

    def test_speech_page_uses_csp_safe_preview_script(self, client, monkeypatch):
        from acb_large_print_web.routes import speech as speech_routes

        monkeypatch.setattr(
            speech_routes,
            "get_engine_status",
            lambda: {
                "kokoro": {
                    "installed": True,
                    "models_present": True,
                    "ready": True,
                    "voices_available": ["af_bella"],
                    "model_dir": "/tmp/speech_models",
                    "setup_commands": [],
                },
                "piper": {
                    "installed": False,
                    "models_present": False,
                    "ready": False,
                    "voices_available": [],
                    "model_dir": "/tmp/speech_models/piper",
                    "setup_commands": [],
                },
            },
        )

        resp = client.get("/speech/")

        assert resp.status_code == 200
        assert b'data-preview-url="/speech/preview"' in resp.data
        assert b"/static/speech.js" in resp.data
        assert b"oninput=" not in resp.data
        assert b"URL.createObjectURL" not in resp.data
        assert re.search(
            rb'<input type="radio" name="voice" value="kokoro:af_bella"\s+checked>',
            resp.data,
        )

    def test_quick_start_choose_shows_speech_action(self, client):
        data = {"file": (io.BytesIO(b"# Hello\n\nThis is speech-ready markdown."), "sample.md")}
        start = client.post("/process/", data=data, content_type="multipart/form-data")
        assert start.status_code == 302
        choose = client.get(start.location)
        assert choose.status_code == 200
        assert b"Speech" in choose.data

    def test_speech_prepare_preview_and_download_document(self, client, monkeypatch):
        from acb_large_print_web.routes import speech as speech_routes

        def _fake_wav(_voice_id, _text, speed=1.0, pitch=0):
            buf = io.BytesIO()
            with wave.open(buf, "wb") as wf:
                wf.setnchannels(1)
                wf.setsampwidth(2)
                wf.setframerate(22050)
                wf.writeframes(b"\0\0" * 2205)
            return buf.getvalue(), "glow-speech-test.wav"

        monkeypatch.setattr(speech_routes, "synthesize", _fake_wav)
        monkeypatch.setattr(speech_routes, "synthesize_document_text", _fake_wav)

        prep = client.post(
            "/speech/prepare",
            data={"document": (io.BytesIO(b"# Title\n\nFirst sentence. Second sentence. Third sentence."), "doc.md")},
            content_type="multipart/form-data",
        )
        assert prep.status_code == 200
        payload = prep.get_json()
        assert payload["ok"] is True
        assert payload["token"]
        assert payload["preview_text"]
        assert payload["estimate_processing_seconds"] > 0

        preview = client.post(
            "/speech/document-preview",
            data={
                "token": payload["token"],
                "voice": "kokoro:af_bella",
                "speed": "1.0",
                "pitch": "0",
            },
        )
        assert preview.status_code == 200
        assert preview.headers.get("Content-Type") == "audio/wav"

        download = client.post(
            "/speech/document-download",
            data={
                "token": payload["token"],
                "voice": "kokoro:af_bella",
                "speed": "1.0",
                "pitch": "0",
            },
        )
        assert download.status_code == 200
        assert download.headers.get("Content-Type") in {"audio/wav", "audio/mpeg"}
        assert "attachment;" in (download.headers.get("Content-Disposition") or "")

    def test_caddy_csp_allows_speech_preview_fetch_and_blob_audio(self):
        caddyfile = Path("web/Caddyfile").read_text(encoding="utf-8")
        caddyfile_example = Path("web/Caddyfile.example").read_text(encoding="utf-8")

        assert "connect-src 'self'" in caddyfile
        assert "connect-src 'self'" in caddyfile_example
        assert "media-src 'self' blob:" in caddyfile
        assert "media-src 'self' blob:" in caddyfile_example

    def test_deploy_seeds_piper_default_voice(self):
        deploy_script = Path("scripts/deploy-app.sh").read_text(encoding="utf-8")
        dockerfile = Path("web/Dockerfile").read_text(encoding="utf-8")

        assert "Ensuring curated Piper voice model files are present" in deploy_script
        assert "en_US-lessac-medium" in deploy_script
        assert "en_US-amy-medium" in deploy_script
        assert "en_GB-alan-medium" in deploy_script
        assert "https://huggingface.co/rhasspy/piper-voices/resolve/main" in deploy_script
        assert "Piper voice seeding complete." in dockerfile


class TestPdfAudit:
    """Upload PDF files and verify audit works."""

    def test_audit_pdf_full(self, client):
        pdf_data = _make_fake_pdf()
        data = {
            "document": (pdf_data, "test.pdf"),
            "mode": "full",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Audit Report" in resp.data or b"Score" in resp.data

    def test_audit_pdf_quick(self, client):
        pdf_data = _make_fake_pdf()
        data = {
            "document": (pdf_data, "test.pdf"),
            "mode": "quick",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_corrupt_pdf_rejected(self, client):
        data = {"document": (io.BytesIO(b"not a pdf"), "test.pdf")}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"does not appear to be a valid" in resp.data

    def test_fix_pdf_returns_warning(self, client):
        """PDF fix should work but return a 'cannot be auto-fixed' warning."""
        pdf_data = _make_fake_pdf()
        data = {"document": (pdf_data, "test.pdf")}
        resp = client.post("/fix/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200


# ============================================================
# Rules help URLs
# ============================================================


class TestRulesHelpUrls:
    """Verify help URL data layer returns valid entries."""

    def test_help_urls_map_not_empty(self):
        from acb_large_print_web.rules import get_help_urls_map

        urls_map = get_help_urls_map()
        assert isinstance(urls_map, dict)
        assert len(urls_map) > 0

    def test_help_urls_contain_valid_dicts(self):
        from acb_large_print_web.rules import get_help_urls_map

        urls_map = get_help_urls_map()
        for rule_id, links in urls_map.items():
            assert isinstance(rule_id, str)
            assert isinstance(links, list)
            for entry in links:
                assert "label" in entry, f"Missing 'label' for {rule_id}"
                assert "url" in entry, f"Missing 'url' for {rule_id}"
                assert entry["url"].startswith(
                    "https://"
                ), f"Bad URL for {rule_id}: {entry['url']}"

    def test_wcag_rules_have_w3c_urls(self):
        from acb_large_print_web.rules import get_help_urls_map

        urls_map = get_help_urls_map()
        # At least some WCAG-linked rules should have w3.org URLs
        w3_count = sum(
            1
            for links in urls_map.values()
            for entry in links
            if "w3.org" in entry["url"]
        )
        assert w3_count > 0, "No W3C URLs found in help URLs map"


# ============================================================
# Upload validation for new formats
# ============================================================


class TestUploadValidation:
    """Verify upload.py accepts/rejects the new file formats."""

    def test_md_extension_accepted(self, client):
        md_data = _make_fake_md()
        data = {"document": (md_data, "test.md"), "mode": "full"}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_pdf_extension_accepted(self, client):
        pdf_data = _make_fake_pdf()
        data = {"document": (pdf_data, "test.pdf"), "mode": "full"}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_home_shows_md_and_pdf_formats(self, client):
        resp = client.get("/")
        assert b"Markdown" in resp.data or b".md" in resp.data
        assert b"PDF" in resp.data or b".pdf" in resp.data
        assert b"ePub" in resp.data or b".epub" in resp.data

    def test_audit_form_accepts_md_and_pdf(self, client):
        resp = client.get("/audit/")
        assert b".md" in resp.data
        assert b".pdf" in resp.data
        assert b".epub" in resp.data


# ============================================================
# ePub audit
# ============================================================


def _make_fake_epub() -> io.BytesIO:
    """Create a minimal valid .epub file (ZIP with OPF metadata)."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("mimetype", "application/epub+zip")
        zf.writestr(
            "META-INF/container.xml",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<container xmlns="urn:oasis:names:tc:opendocument:xmlns:container" version="1.0">'
            '<rootfiles><rootfile full-path="content.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles>'
            "</container>",
        )
        # OPF with no title and no language -- should trigger EPUB-TITLE + EPUB-LANGUAGE
        zf.writestr(
            "content.opf",
            '<?xml version="1.0" encoding="UTF-8"?>'
            '<package xmlns="http://www.idpf.org/2007/opf" version="3.0">'
            '<metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            "</metadata>"
            "<manifest>"
            '<item id="ch1" href="ch1.xhtml" media-type="application/xhtml+xml"/>'
            "</manifest>"
            '<spine><itemref idref="ch1"/></spine>'
            "</package>",
        )
        zf.writestr(
            "ch1.xhtml",
            "<html><body><h1>Chapter 1</h1><p>Content</p></body></html>",
        )
    buf.seek(0)
    return buf


class TestEpubAudit:
    """Upload ePub files and verify audit works."""

    def test_audit_epub_full(self, client):
        epub_data = _make_fake_epub()
        data = {
            "document": (epub_data, "test.epub"),
            "mode": "full",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert b"Audit Report" in resp.data or b"Score" in resp.data

    def test_audit_epub_quick(self, client):
        epub_data = _make_fake_epub()
        data = {
            "document": (epub_data, "test.epub"),
            "mode": "quick",
        }
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_corrupt_epub_rejected(self, client):
        data = {"document": (io.BytesIO(b"not a zip archive"), "test.epub")}
        resp = client.post("/audit/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"does not appear to be a valid" in resp.data

    def test_fix_epub_returns_warning(self, client):
        """ePub fix should work but return audit-only guidance."""
        epub_data = _make_fake_epub()
        data = {"document": (epub_data, "test.epub")}
        resp = client.post("/fix/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200

    def test_epub_auditor_finds_missing_metadata(self):
        """Direct unit test: the epub auditor detects missing title and language."""
        import tempfile
        from acb_large_print.epub_auditor import audit_epub

        epub_data = _make_fake_epub()
        with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as f:
            f.write(epub_data.read())
            f.flush()
            result = audit_epub(f.name)
        rule_ids = {finding.rule_id for finding in result.findings}
        assert "EPUB-TITLE" in rule_ids, f"Expected EPUB-TITLE, got {rule_ids}"
        assert "EPUB-LANGUAGE" in rule_ids, f"Expected EPUB-LANGUAGE, got {rule_ids}"
        assert (
            "EPUB-NAV-DOCUMENT" in rule_ids
        ), f"Expected EPUB-NAV-DOCUMENT, got {rule_ids}"
        assert "EPUB-ACCESSIBILITY-METADATA" in rule_ids


# ============================================================
# DAISY Pipeline conversion
# ============================================================


class TestPipelineConversion:
    """Tests for DAISY Pipeline integration on the /convert/ route."""

    def test_convert_form_shows_pipeline_when_available(self, client, monkeypatch):
        """When Pipeline is available, the form should show the EPUB/DAISY option."""
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.pipeline_available", lambda: True
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.get_available_conversions",
            lambda: {
                "html-to-epub": {
                    "label": "HTML to EPUB 3",
                    "script": "html-to-epub3",
                    "description": "Packages HTML into an accessible EPUB 3",
                    "input_ext": ".html",
                    "output_ext": ".epub",
                },
            },
        )
        resp = client.get("/convert/")
        assert resp.status_code == 200
        assert b"EPUB or DAISY format" in resp.data
        assert b"DAISY Pipeline" in resp.data

    def test_convert_form_hides_pipeline_when_unavailable(self, client, monkeypatch):
        """When Pipeline is not installed, the EPUB/DAISY radio should not appear."""
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.pipeline_available", lambda: False
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.get_available_conversions",
            lambda: {},
        )
        resp = client.get("/convert/")
        assert resp.status_code == 200
        assert b'value="to-pipeline"' not in resp.data

    def test_pipeline_not_installed_returns_error(self, client, monkeypatch):
        """Submitting a Pipeline conversion when Pipeline is absent returns an error."""
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.pipeline_available", lambda: False
        )
        html_content = b"<html><body><h1>Test</h1></body></html>"
        data = {
            "document": (io.BytesIO(html_content), "test.html"),
            "direction": "to-pipeline",
            "pipeline_conversion": "html-to-epub",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"Pipeline" in resp.data

    def test_pipeline_invalid_conversion_key(self, client, monkeypatch):
        """Submitting an unknown Pipeline conversion key returns an error."""
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.pipeline_available", lambda: True
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.get_available_conversions",
            lambda: {},
        )
        html_content = b"<html><body><h1>Test</h1></body></html>"
        data = {
            "document": (io.BytesIO(html_content), "test.html"),
            "direction": "to-pipeline",
            "pipeline_conversion": "nonexistent-conversion",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 400
        assert b"not available" in resp.data

    def test_pipeline_html_to_epub_success(self, client, monkeypatch, tmp_path):
        """Mock a successful Pipeline HTML-to-EPUB conversion."""
        # Create a fake EPUB output file
        fake_epub = tmp_path / "output.epub"
        fake_epub.write_bytes(b"PK\x03\x04fake-epub-content")

        def mock_convert(src_path, conversion_key, output_dir=None):
            return fake_epub, "Converted to output.epub"

        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.pipeline_available", lambda: True
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.get_available_conversions",
            lambda: {
                "html-to-epub": {
                    "label": "HTML to EPUB 3",
                    "script": "html-to-epub3",
                    "description": "Test",
                    "input_ext": ".html",
                    "output_ext": ".epub",
                },
            },
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.convert_with_pipeline",
            mock_convert,
        )
        html_content = b"<html><body><h1>Test</h1></body></html>"
        data = {
            "document": (io.BytesIO(html_content), "test.html"),
            "direction": "to-pipeline",
            "pipeline_conversion": "html-to-epub",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert resp.content_type == "application/epub+zip"
        assert b"output.epub" in resp.headers.get("Content-Disposition", "").encode()

    def test_pipeline_epub_to_daisy_returns_zip(self, client, monkeypatch, tmp_path):
        """Mock a Pipeline EPUB-to-DAISY conversion that returns a directory (zipped)."""
        # Create a fake DAISY output directory
        daisy_dir = tmp_path / "daisy_output"
        daisy_dir.mkdir()
        (daisy_dir / "ncc.html").write_text("<html><body>NCC</body></html>")
        (daisy_dir / "content.smil").write_text("<smil>test</smil>")

        def mock_convert(src_path, conversion_key, output_dir=None):
            return daisy_dir, "Converted to DAISY 2.02"

        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.pipeline_available", lambda: True
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.get_available_conversions",
            lambda: {
                "epub-to-daisy202": {
                    "label": "EPUB to DAISY 2.02",
                    "script": "epub3-to-daisy202",
                    "description": "Test",
                    "input_ext": ".epub",
                    "output_ext": "",
                },
            },
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.convert_with_pipeline",
            mock_convert,
        )
        epub_data = _make_fake_epub()
        data = {
            "document": (epub_data, "test.epub"),
            "direction": "to-pipeline",
            "pipeline_conversion": "epub-to-daisy202",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert resp.content_type == "application/zip"
        assert b".zip" in resp.headers.get("Content-Disposition", "").encode()

    def test_pipeline_legacy_direction_format(self, client, monkeypatch, tmp_path):
        """The route also accepts pipeline-{key} as the direction value (legacy)."""
        fake_epub = tmp_path / "output.epub"
        fake_epub.write_bytes(b"PK\x03\x04fake-epub")

        def mock_convert(src_path, conversion_key, output_dir=None):
            return fake_epub, "Converted"

        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.pipeline_available", lambda: True
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.get_available_conversions",
            lambda: {
                "html-to-epub": {
                    "label": "HTML to EPUB 3",
                    "script": "html-to-epub3",
                    "description": "Test",
                    "input_ext": ".html",
                    "output_ext": ".epub",
                },
            },
        )
        monkeypatch.setattr(
            "acb_large_print_web.routes.convert.convert_with_pipeline",
            mock_convert,
        )
        html_content = b"<html><body><h1>Test</h1></body></html>"
        data = {
            "document": (io.BytesIO(html_content), "test.html"),
            "direction": "pipeline-html-to-epub",
        }
        resp = client.post("/convert/", data=data, content_type="multipart/form-data")
        assert resp.status_code == 200
        assert resp.content_type == "application/epub+zip"


class TestPipelineConverterUnit:
    """Unit tests for pipeline_converter.py functions (no dp2 required)."""

    def test_pipeline_conversions_dict_complete(self):
        from acb_large_print.pipeline_converter import PIPELINE_CONVERSIONS

        assert "docx-to-epub" in PIPELINE_CONVERSIONS
        assert "html-to-epub" in PIPELINE_CONVERSIONS
        assert "epub-to-daisy202" in PIPELINE_CONVERSIONS
        assert "epub-to-daisy3" in PIPELINE_CONVERSIONS

    def test_pipeline_input_extensions(self):
        from acb_large_print.pipeline_converter import PIPELINE_INPUT_EXTENSIONS

        assert ".docx" in PIPELINE_INPUT_EXTENSIONS
        assert ".html" in PIPELINE_INPUT_EXTENSIONS
        assert ".epub" in PIPELINE_INPUT_EXTENSIONS

    def test_pipeline_available_returns_bool(self):
        from acb_large_print.pipeline_converter import pipeline_available

        result = pipeline_available()
        assert isinstance(result, bool)

    def test_convert_with_pipeline_missing_file(self):
        from acb_large_print.pipeline_converter import convert_with_pipeline

        with pytest.raises(FileNotFoundError):
            convert_with_pipeline(Path("/nonexistent/file.html"), "html-to-epub")

    def test_convert_with_pipeline_bad_key(self, tmp_path):
        from acb_large_print.pipeline_converter import convert_with_pipeline

        f = tmp_path / "test.html"
        f.write_text("<html></html>")
        with pytest.raises(ValueError, match="Unknown conversion"):
            convert_with_pipeline(f, "nonexistent-key")

    def test_convert_with_pipeline_wrong_extension(self, tmp_path):
        from acb_large_print.pipeline_converter import convert_with_pipeline

        f = tmp_path / "test.txt"
        f.write_text("not html")
        with pytest.raises(ValueError, match="expects .html"):
            convert_with_pipeline(f, "html-to-epub")

    def test_convert_with_pipeline_not_installed(self, tmp_path, monkeypatch):
        from acb_large_print import pipeline_converter

        monkeypatch.setattr(pipeline_converter, "pipeline_available", lambda: False)
        f = tmp_path / "test.html"
        f.write_text("<html></html>")
        with pytest.raises(RuntimeError, match="not reachable"):
            pipeline_converter.convert_with_pipeline(f, "html-to-epub")
