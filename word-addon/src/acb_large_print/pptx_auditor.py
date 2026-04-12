"""Audit PowerPoint presentations for accessibility compliance."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from . import constants as C
from .auditor import AuditResult

# Non-descriptive link text
_BAD_LINK_RE = re.compile(
    r"^(click here|here|link|more|read more|learn more|download|info)$",
    re.IGNORECASE,
)
_URL_RE = re.compile(r"^https?://", re.IGNORECASE)

# Fake bullet characters
_FAKE_BULLET_RE = re.compile(r"^[\u2022\u2023\u25CB\u25CF\u25E6\u2043\u2219\-\*]\s")


def audit_presentation(file_path: str | Path) -> AuditResult:
    """Audit a PowerPoint presentation for accessibility issues.

    Returns an AuditResult populated with findings scoped to PPTX rules.
    """
    file_path = Path(file_path)
    result = AuditResult(file_path=str(file_path))

    prs = Presentation(str(file_path))

    # -- Document properties --
    if not (prs.core_properties.title and prs.core_properties.title.strip()):
        result.add("PPTX-TITLE", "Presentation title is not set in document properties.")

    if not (prs.core_properties.author and prs.core_properties.author.strip()):
        result.add("ACB-DOC-AUTHOR", "Presentation author is not set in document properties.")

    slide_titles: list[str] = []
    slides_without_notes = 0
    slides_with_visuals_no_notes = 0

    for idx, slide in enumerate(prs.slides, 1):
        loc = f"Slide {idx}"

        # -- Slide title --
        title_shape = slide.shapes.title
        title_text = ""
        if title_shape is not None:
            title_text = (title_shape.text or "").strip()

        if not title_text:
            result.add(
                "PPTX-SLIDE-TITLE",
                "Slide has no title or the title placeholder is empty.",
                loc,
            )
        else:
            slide_titles.append(title_text)

        # -- Reading order --
        _check_reading_order(slide, result, loc, title_shape)

        # -- Images / shapes alt text --
        has_visual = False
        for shape in slide.shapes:
            if shape.shape_type in (
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.LINKED_PICTURE,
            ):
                has_visual = True
                _check_shape_alt_text(shape, result, loc)
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                has_visual = True
                _check_chart_alt_text(shape, result, loc)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                has_visual = True
                _check_shape_alt_text(shape, result, loc)
            elif hasattr(shape, "image"):
                has_visual = True
                _check_shape_alt_text(shape, result, loc)

        # -- Text checks (font size, fake lists, links) --
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                _check_paragraph(para, shape, result, loc)

        # -- Speaker notes --
        notes = slide.notes_slide if slide.has_notes_slide else None
        has_notes = notes and notes.notes_text_frame and notes.notes_text_frame.text.strip()
        if not has_notes:
            slides_without_notes += 1
            if has_visual:
                slides_with_visuals_no_notes += 1

    # -- Duplicate slide titles --
    _check_duplicate_titles(slide_titles, result)

    # -- Speaker notes summary --
    if slides_with_visuals_no_notes > 0:
        result.add(
            "PPTX-SPEAKER-NOTES",
            f"{slides_with_visuals_no_notes} slide(s) with visual content but no speaker notes.",
        )

    # Sort findings by severity
    severity_order = {
        C.Severity.CRITICAL: 0,
        C.Severity.HIGH: 1,
        C.Severity.MEDIUM: 2,
        C.Severity.LOW: 3,
    }
    result.findings.sort(key=lambda f: severity_order.get(f.severity, 4))

    return result


# ---------------------------------------------------------------------------
# Individual check helpers
# ---------------------------------------------------------------------------

def _check_reading_order(slide, result: AuditResult, loc: str, title_shape) -> None:
    """Check that reading order roughly matches visual order."""
    shapes = [(s.left or 0, s.top or 0, s.name, s) for s in slide.shapes]
    if len(shapes) < 2:
        return

    # Visual order: top-to-bottom, then left-to-right
    visual = sorted(shapes, key=lambda s: (s[1], s[0]))
    xml_order = [s[2] for s in shapes]
    visual_order = [s[2] for s in visual]

    if xml_order != visual_order:
        result.add(
            "PPTX-READING-ORDER",
            "Reading order differs from visual layout. Check Selection Pane.",
            loc,
        )

    # Title should be first in reading order
    if title_shape is not None and shapes:
        # In PowerPoint, screen readers traverse back-to-front in the XML
        first_in_xml = shapes[0][2] if shapes else ""
        if title_shape.name and title_shape.name != first_in_xml:
            # Only flag if title isn't at the start
            xml_names = [s[2] for s in shapes]
            if title_shape.name in xml_names:
                title_idx = xml_names.index(title_shape.name)
                if title_idx > 0:
                    result.add(
                        "PPTX-TITLE-READING-ORDER",
                        "Title placeholder is not first in reading order.",
                        loc,
                    )


def _check_shape_alt_text(shape, result: AuditResult, loc: str) -> None:
    """Check a shape for alternative text."""
    # python-pptx stores alt text via the XML descr attribute
    desc = ""
    if hasattr(shape, "_element"):
        # Try nvSpPr/cNvPr or nvPicPr/cNvPr
        for tag in ("p:nvSpPr", "p:nvPicPr", "p:nvGrpSpPr"):
            nvPr = shape._element.find(qn(tag))
            if nvPr is not None:
                cNvPr = nvPr.find(qn("p:cNvPr"))
                if cNvPr is not None:
                    desc = cNvPr.get("descr", "") or ""
                    break
        if not desc:
            desc = shape._element.get("descr", "") or ""

    if not desc.strip():
        result.add(
            "ACB-MISSING-ALT-TEXT",
            f"Shape '{shape.name}' has no alternative text.",
            loc,
        )


def _check_chart_alt_text(shape, result: AuditResult, loc: str) -> None:
    """Check a chart shape for alt text or title."""
    desc = ""
    if hasattr(shape, "_element"):
        desc = shape._element.get("descr", "") or ""
    if not desc.strip():
        result.add(
            "PPTX-CHART-ALT-TEXT",
            f"Chart '{shape.name}' has no alt text describing its key finding.",
            loc,
        )


def _check_paragraph(para, shape, result: AuditResult, loc: str) -> None:
    """Check a text paragraph for font size, fake lists, and links."""
    text = para.text or ""

    # Font size check
    for run in para.runs:
        if run.font and run.font.size:
            size_pt = run.font.size.pt
            if size_pt < 18:
                result.add(
                    "PPTX-SMALL-FONT",
                    f"Text '{text[:40]}...' is {size_pt}pt (minimum 18pt recommended).",
                    loc,
                )
                break  # One finding per paragraph

    # Fake bullet
    if _FAKE_BULLET_RE.match(text):
        result.add(
            "ACB-FAKE-LIST",
            f"Manually typed bullet character in '{text[:40]}'.",
            loc,
        )

    # Hyperlinks
    for run in para.runs:
        if run.hyperlink and run.hyperlink.address:
            link_text = run.text.strip()
            if _BAD_LINK_RE.match(link_text):
                result.add(
                    "ACB-LINK-TEXT",
                    f"Non-descriptive link text '{link_text}'.",
                    loc,
                )
            elif _URL_RE.match(link_text):
                result.add(
                    "ACB-LINK-TEXT",
                    f"Raw URL used as link text: '{link_text[:60]}'.",
                    loc,
                )


def _check_duplicate_titles(titles: list[str], result: AuditResult) -> None:
    """Check for duplicate slide titles."""
    seen: dict[str, int] = {}
    for title in titles:
        lower = title.lower()
        seen[lower] = seen.get(lower, 0) + 1

    for title_lower, count in seen.items():
        if count > 1:
            result.add(
                "PPTX-DUPLICATE-SLIDE-TITLE",
                f"Title '{title_lower}' is used on {count} slides.",
            )
